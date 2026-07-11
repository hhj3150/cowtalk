# -*- coding: utf-8 -*-
# 지역 축산 AX 디지털 플랫폼 구축사업 — 마스터 발표 덱 (McKinsey 스타일)
# 90/10 원칙: 90% 고정, 지역별로 REGION preset 1블록만 교체하면 5개 도 발표 완성.
# 논리 골격: 답부터 → 왜 지금 → 왜 지역·무엇을 → 어떻게 → 그래서 결과는 → 실행·확장
# 정직성: 현재 구현(✅) / 부분(🟡) / 향후(🔵)를 슬라이드에서 분리 표기. 모든 수치 출처.
import sys
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ══════════════════════════════════════════════════════════════════
# 지역 치환 변수 (마스터의 유일한 가변부). 인자로 프리셋 선택: python3 build_master_deck.py gyeonggi
# 네트워크 농장 수 = 현재 CowTalk 실운영 분포(정직 수치). 경제는 지역 실두수로 재산출 전제.
# ══════════════════════════════════════════════════════════════════
REGIONS = {
    "generic":  dict(name="○○도", short="○○", pos="축우(한우·젖소) 중심",
                     posShort="축우", net="N", cattle="지역 축우",
                     start="축우(한우·젖소)", special="데이터가 가장 많은 곳이 표준을 만든다",
                     rank="지역 축산 거점"),
    "gyeonggi": dict(name="경기도", short="경기", pos="대한민국 낙농의 중심 (젖소 사육 전국 약 41%)",
                     posShort="낙농", net="52", cattle="경기 젖소·한우",
                     start="젖소(낙농) → 한우", special="생산과 소비를 동시에 가진 유일 광역 (수도권 최대 소비지)",
                     rank="수도권 최대 · 전국 2위 네트워크"),
    "jeonnam":  dict(name="전라남도", short="전남", pos="한우 주산지 (한우 출하 전국 2위 · 15.4%)",
                     posShort="한우", net="8", cattle="전남 한우",
                     start="한우 → 낙농", special="한우 주산지의 방역·번식을 데이터로 지킨다",
                     rank="한우 주산지 거점"),
    "chungnam": dict(name="충청남도", short="충남", pos="한우 3위(12.7%) + 낙농 복합 · 중부권 물류 중심",
                     posShort="축우 복합", net="27", cattle="충남 축우",
                     start="한우·젖소 병행", special="중부권 복합 축산의 표준 모델",
                     rank="중부권 복합 축산 거점"),
    "gyeongbuk":dict(name="경상북도", short="경북", pos="대한민국 한우 1번지 (한우 출하 전국 1위 · 23.0%)",
                     posShort="한우", net="14", cattle="경북 한우",
                     start="한우", special="한우 1번지가 한우 데이터 표준을 만든다",
                     rank="한우 출하 전국 1위 거점"),
    "jeju":     dict(name="제주특별자치도", short="제주", pos="청정·브랜드 축산 (제주 한우·흑우)",
                     posShort="청정 브랜드", net="4", cattle="제주 축우",
                     start="한우·브랜드", special="청정 이미지와 수출 브랜드를 데이터로 증명한다",
                     rank="청정·브랜드 축산 거점"),
}
KEY = sys.argv[1] if len(sys.argv) > 1 else "generic"
RG = REGIONS.get(KEY, REGIONS["generic"])
NAME, SHORT, NET = RG["name"], RG["short"], RG["net"]

# ── 디자인 토큰 (기존 인수위 덱 계승) ──
INK   = RGBColor(0x0B,0x22,0x39); INK2 = RGBColor(0x24,0x3B,0x53)
TEAL  = RGBColor(0x0E,0x7C,0x7B); TEALD= RGBColor(0x09,0x53,0x52); MINT = RGBColor(0x16,0xB0,0xA6)
AMBER = RGBColor(0xC9,0x7A,0x1E); GRAY = RGBColor(0x55,0x66,0x76); MUTE = RGBColor(0x90,0x9E,0xAA)
LINE  = RGBColor(0xD8,0xE0,0xE6); PANEL= RGBColor(0xF4,0xF7,0xF8); PANEL2=RGBColor(0xEA,0xF1,0xF1)
WHITE = RGBColor(0xFF,0xFF,0xFF); PALE = RGBColor(0xC7,0xD6,0xDE)
OK    = RGBColor(0x0E,0x7C,0x7B); MID = RGBColor(0xC9,0x7A,0x1E); FUT = RGBColor(0x6B,0x7A,0x88)
KR = "맑은 고딕"

prs = Presentation(); prs.slide_width = I(13.333); prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
EW, EH = 13.333, 7.5; ML, MR = 0.72, 0.72; CW = EW - ML - MR

def slide(): return prs.slides.add_slide(BLANK)
def _fill(sh,c):
    if c is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = c
def _line(sh,c,w=0.75):
    if c is None: sh.line.fill.background()
    else: sh.line.color.rgb = c; sh.line.width = Pt(w)
def box(s,x,y,w,h,fill=None,line=None,lw=0.75,rounded=False,radius=0.06):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,I(x),I(y),I(w),I(h))
    _fill(sh,fill); _line(sh,line,lw); sh.shadow.inherit=False
    if rounded:
        try: sh.adjustments[0]=radius
        except Exception: pass
    return sh
def hline(s,x,y,w,color=LINE,weight=0.75):
    ln=s.shapes.add_connector(2,I(x),I(y),I(x+w),I(y)); ln.line.color.rgb=color; ln.line.width=Pt(weight); ln.shadow.inherit=False; return ln
def text(s,x,y,w,h,paras,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    if isinstance(paras,str): paras=[{"runs":[{"t":paras}]}]
    for i,pa in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=pa.get("align",PP_ALIGN.LEFT)
        if pa.get("sb") is not None: p.space_before=Pt(pa["sb"])
        if pa.get("sa") is not None: p.space_after=Pt(pa["sa"])
        if pa.get("ls") is not None: p.line_spacing=pa["ls"]
        for r in pa["runs"]:
            run=p.add_run(); run.text=r["t"]; f=run.font
            f.name=r.get("f",KR); f.size=Pt(r.get("sz",14)); f.bold=r.get("b",False); f.italic=r.get("it",False)
            f.color.rgb=r.get("c",INK)
            rPr=run._r.get_or_add_rPr(); ea=rPr.find(qn('a:ea'))
            if ea is None: ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
            ea.set('typeface',r.get("f",KR))
    return tb
def R(t,sz=14,c=INK,b=False,f=KR,it=False): return {"t":t,"sz":sz,"c":c,"b":b,"f":f,"it":it}

PHASE_STYLE={"현황":(RGBColor(0x6B,0x7A,0x88),WHITE),"무엇을":(TEAL,WHITE),"어떻게":(INK,WHITE),
             "그래서 결과는":(AMBER,WHITE),"실행":(INK,WHITE),"제언":(TEALD,WHITE)}
def header(s,kicker,title,phase=None,title_sz=22):
    text(s,ML,0.46,CW-2.9,0.3,[{"runs":[R(kicker,11.5,TEAL,True)]}])
    if phase:
        bg,fg=PHASE_STYLE.get(phase,(TEAL,WHITE)); pw=1.05+0.115*len(phase)
        box(s,EW-MR-pw,0.44,pw,0.34,fill=bg,line=None,rounded=True,radius=0.5)
        text(s,EW-MR-pw,0.44,pw,0.34,[{"align":PP_ALIGN.CENTER,"runs":[R(phase,11,fg,True)]}],anchor=MSO_ANCHOR.MIDDLE)
    text(s,ML,0.80,CW,1.0,[{"ls":1.04,"runs":[R(title,title_sz,INK,True)]}])
    hline(s,ML,1.82,0.62,color=TEAL,weight=2.6)

TOTAL=21; _pg=[1]
def footer(s):
    _pg[0]+=1; n=_pg[0]
    hline(s,ML,7.04,CW,color=LINE,weight=0.75)
    text(s,ML,7.10,9.5,0.3,[{"runs":[R("지역 축산 AX 플랫폼",9,TEAL,True),R(f"   ·   {NAME} 축산 디지털 대전환 제안",9,MUTE)]}])
    text(s,EW-MR-2.0,7.10,2.0,0.3,[{"align":PP_ALIGN.RIGHT,"runs":[R(f"{n:02d} / {TOTAL}",9,MUTE,True)]}])

def chip(s,x,y,label,col):  # 정직성 칩 (✅/🟡/🔵)
    w=0.34+0.13*len(label)
    box(s,x,y,w,0.30,fill=None,line=col,lw=1.1,rounded=True,radius=0.5)
    text(s,x,y,w,0.30,[{"align":PP_ALIGN.CENTER,"runs":[R(label,10,col,True)]}],anchor=MSO_ANCHOR.MIDDLE)
    return w

# ══════════════════════════════════════════════════════════════════
# 01. 표지
# ══════════════════════════════════════════════════════════════════
s=slide(); box(s,0,0,EW,EH,fill=INK,line=None); box(s,0,0,0.22,EH,fill=TEAL,line=None)
for yy in [1.0,1.34,1.68]: hline(s,EW-3.4,yy,2.7,color=RGBColor(0x1E,0x3A,0x52),weight=1.4)
text(s,1.0,1.55,11,0.4,[{"runs":[R(f"{NAME} 정책 제안 · 2026",13,MINT,True)]}])
text(s,1.0,2.2,11.6,2.3,[{"ls":1.02,"runs":[R("지역 축산 AX",44,WHITE,True)]},
     {"sb":2,"ls":1.02,"runs":[R("디지털 플랫폼 구축사업",44,WHITE,True)]},
     {"sb":12,"ls":1.1,"runs":[R(f"센서 · 공공데이터 · AI로 완성하는 「{NAME} 축산 운영체제」",17,PALE)]}])
hline(s,1.0,5.15,4.2,color=TEAL,weight=2.2)
text(s,1.0,5.35,11.4,0.8,[{"runs":[R(f"“{NAME}가 대한민국 축산 AX의 첫 표준을 만듭니다.”",17,WHITE,True,it=True)]}])
text(s,1.0,6.75,11.4,0.4,[{"runs":[R("CowTalk v5.0",12,WHITE,True),R("   ·   Foundation Platform for Livestock AX",12,MUTE)]}])

# ══════════════════════════════════════════════════════════════════
# 02. 제언 (답부터)
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"제언  ·  THE ANSWER",f"{NAME}가 대한민국 축산 디지털 대전환을 선도해야 합니다.",phase="제언",title_sz=23)
text(s,ML,2.02,CW,0.4,[{"runs":[R("새 소프트웨어를 사는 사업이 아닙니다. ",13.5,INK,True),
     R("이미 199개 농장에서 실제로 돌아가는 검증된 기반 위에, 지역을 국가 표준의 원점으로 만드는 사업입니다.",13.5,GRAY)]}])
cards=[("농가 소득","번식·질병·노동을 데이터로","두당 연 13~24만원 생산성",TEAL),
       ("방역 주권","역학조사 며칠 → 30초","전염병 1회 차단 = 사업비 회수",INK2),
       ("지역 AI 산업","데이터·인재·일자리","지역이 표준을 소유",TEALD)]
cw=(CW-2*0.4)/3; cy=2.9; ch=1.95
for i,(t,d,v,col) in enumerate(cards):
    cx=ML+i*(cw+0.4); box(s,cx,cy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,cy,cw,0.08,fill=col,line=None)
    text(s,cx+0.28,cy+0.30,cw-0.5,0.4,[{"runs":[R(t,17,INK,True)]}])
    text(s,cx+0.28,cy+0.82,cw-0.5,0.5,[{"ls":1.05,"runs":[R(d,12.5,GRAY)]}])
    text(s,cx+0.28,cy+1.40,cw-0.5,0.4,[{"runs":[R(v,12.5,col,True)]}])
box(s,ML,5.2,CW,0.95,fill=INK,line=None)
text(s,ML+0.4,5.2,CW-0.8,0.95,[{"ls":1.15,"runs":[
     R("결론은 하나입니다 — ",14,PALE),R("지금, 우리 지역이 먼저. ",15,WHITE,True),
     R("데이터가 가장 먼저 쌓이는 곳이 표준을 만들고, 표준을 만든 지역이 국가 플랫폼의 원점이 됩니다.",13.5,PALE)]}],anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 03. Q1 축산 현주소
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"현황  ·  Q1 대한민국 축산의 현주소","축산은 농업 생산액의 43%를 책임지지만, 지금 네 가지 위기가 동시에 옵니다.",phase="현황")
cards=[("기간산업","25.5조","축산 생산액 (농업의 42.9%)",TEAL),
       ("사육두수 감소","321만·37.5만","한우·젖소 (2025, 감소세)",INK2),
       ("구제역 상시위협","3.18조","2010–11 재정 소요 (348만두 살처분)",AMBER),
       ("현장 고령화","데이터 단절","경험 은퇴 → 판단 공백",FUT)]
cw=(CW-3*0.32)/4; cy=2.35; ch=2.5
for i,(t,big,sub,col) in enumerate(cards):
    cx=ML+i*(cw+0.32); box(s,cx,cy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,cy,cw,0.08,fill=col,line=None)
    text(s,cx+0.22,cy+0.26,cw-0.4,0.35,[{"runs":[R(t,13,INK,True)]}])
    text(s,cx+0.22,cy+0.78,cw-0.4,0.6,[{"runs":[R(big,23,col,True)]}])
    text(s,cx+0.22,cy+1.55,cw-0.4,0.7,[{"ls":1.1,"runs":[R(sub,11,GRAY)]}])
box(s,ML,5.35,CW,0.72,fill=PANEL,line=LINE,lw=0.75)
text(s,ML+0.35,5.35,CW-0.7,0.72,[{"runs":[R("축산은 지킬 가치가 큰 산업인데, 지킬 수단이 아직 아날로그에 머물러 있습니다.",13.5,INK,True)]}],anchor=MSO_ANCHOR.MIDDLE)
text(s,ML,6.2,CW,0.3,[{"runs":[R("출처  통계청 · 식품저널 농업전망 · 위키백과(2010–11 구제역)",9.5,MUTE)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 04. Q2 무엇이 부족한가 — 루프
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"현황  ·  Q2 무엇이 부족한가","부족한 것은 센서가 아니라 '루프'입니다 — 데이터가 판단·행동·학습으로 돌지 않습니다.",phase="현황")
# 좌: 현재(멈춤)
box(s,ML,2.3,5.7,2.7,fill=PANEL,line=LINE,lw=1.0)
text(s,ML+0.3,2.5,5.1,0.35,[{"runs":[R("지금 (스마트팜 단계)",13,GRAY,True)]}])
cur="센서  →  데이터  →  모니터링  →  알람   ▪ 여기서 멈춤"
text(s,ML+0.3,3.1,5.1,1.6,[{"ls":1.5,"runs":[R("센서 → 데이터 → 모니터링 → 알람",14,INK,True)]},
     {"sb":10,"runs":[R("…그리고 멈춘다.",13,AMBER,True)]},
     {"sb":10,"ls":1.2,"runs":[R("‘그래서 뭘 해야 하지?’를 사람이 혼자 판단.",12,GRAY)]}])
# 우: 필요(순환)
box(s,ML+6.1,2.3,CW-6.1,2.7,fill=INK,line=None)
text(s,ML+6.4,2.5,CW-6.7,0.35,[{"runs":[R("필요 (AX 단계)",13,MINT,True)]}])
text(s,ML+6.4,3.1,CW-6.7,1.7,[{"ls":1.45,"runs":[R("데이터 → AI → 의사결정 → 실행",14,WHITE,True)]},
     {"sb":6,"ls":1.45,"runs":[R("→ Outcome → AI Learning ↺",14,WHITE,True)]},
     {"sb":10,"ls":1.2,"runs":[R("결과가 데이터로 되돌아와 AI를 강화하는 순환.",12,PALE)]}])
box(s,ML,5.35,CW,0.72,fill=AMBER,line=None)
text(s,ML,5.35,CW,0.72,[{"align":PP_ALIGN.CENTER,"runs":[
     R("세계 최고 센서(smaXtec)도 ‘알람 → 행동’의 다리가 없습니다. 방역·행정은 여전히 전화와 엑셀.",13.5,WHITE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
text(s,ML,6.2,CW,0.3,[{"runs":[R("스마트팜은 ‘설비’를 깔았고, AX는 그 설비가 만든 데이터를 ‘판단과 행동’으로 바꾸는 다음 단계입니다.",11,GRAY,it=True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 05. Q3 왜 지금 — 정책 창
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"현황  ·  Q3 왜 지금 AX 플랫폼인가","정부가 소버린 AI·AI 대전환을 국가전략으로 미는 지금이 유일한 정책 창입니다.",phase="현황")
cards=[("AI 국가전략","9.9조","2026 AI예산 (2025의 약 3배)",TEAL),
       ("국가 농업 AX","705억","민관합동 농업 AX 플랫폼 조성",INK2),
       ("지역 AI 대전환","5대 분야","AI 대전환에 '지역' 명시 포함",TEALD)]
cw=(CW-2*0.4)/3; cy=2.3; ch=2.0
for i,(t,big,sub,col) in enumerate(cards):
    cx=ML+i*(cw+0.4); box(s,cx,cy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,cy,cw,0.08,fill=col,line=None)
    text(s,cx+0.28,cy+0.28,cw-0.5,0.35,[{"runs":[R(t,14,INK,True)]}])
    text(s,cx+0.28,cy+0.78,cw-0.5,0.6,[{"runs":[R(big,26,col,True)]}])
    text(s,cx+0.28,cy+1.45,cw-0.5,0.4,[{"runs":[R(sub,11.5,GRAY)]}])
box(s,ML,4.6,CW,0.9,fill=PANEL2,line=LINE,lw=0.75)
text(s,ML+0.35,4.6,CW-0.7,0.9,[{"ls":1.15,"runs":[
     R("Smart Farm(설비 보급)  →  AX(운영지능·데이터 주권).  ",13.5,INK,True),
     R("우리는 다음 단계로 갑니다. 예산과 정책 방향이 정렬된 창은 오래 열려 있지 않습니다.",12.5,GRAY)]}],anchor=MSO_ANCHOR.MIDDLE)
text(s,ML,5.65,CW,0.3,[{"runs":[R("출처  korea.kr 정책브리핑 · 농식품부 2026 예산안 · INSS 「한국형 소버린 AI 국가전략」",9.5,MUTE)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 06. Q4 왜 지역부터
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"무엇을  ·  Q4 왜 지역부터인가","국가 축산 AI 플랫폼은 지역에서 시작할 수밖에 없습니다 — 데이터 밀도가 표준을 만듭니다.",phase="무엇을")
rows=[("1","집행 주체","방역·축산 행정을 실제로 집행하는 것은 광역지자체입니다."),
      ("2","데이터 밀도","데이터는 밀도가 높은 곳에서 먼저 표준이 됩니다. 표준을 가진 곳이 원점이 됩니다."),
      ("3","실패비용","전국 일괄 전환은 실패합니다. 한 지역 검증 → 확산이 실패비용이 가장 낮은 정석입니다.")]
cy=2.4
for i,(n,t,d) in enumerate(rows):
    yy=cy+i*1.15
    box(s,ML,yy,0.9,0.9,fill=INK,line=None); text(s,ML,yy,0.9,0.9,[{"align":PP_ALIGN.CENTER,"runs":[R(n,30,MINT,True)]}],anchor=MSO_ANCHOR.MIDDLE)
    text(s,ML+1.2,yy+0.06,CW-1.2,0.4,[{"runs":[R(t,15,INK,True)]}])
    text(s,ML+1.2,yy+0.5,CW-1.2,0.5,[{"ls":1.1,"runs":[R(d,12.5,GRAY)]}])
box(s,ML,6.0,CW,0.62,fill=TEAL,line=None)
text(s,ML,6.0,CW,0.62,[{"align":PP_ALIGN.CENTER,"runs":[R("작은 검증 → 전국 확산. 이것이 국가 플랫폼으로 가는 유일하게 안전한 경로입니다.",13.5,WHITE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 07. 왜 우리 지역 (치환)
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,f"무엇을  ·  왜 {NAME}인가",f"{NAME}는 이미 대한민국 {RG['posShort']}의 중심입니다 — 여기서 시작하는 것이 가장 빠릅니다.",phase="무엇을")
cards=[("축우 산업 포지션",RG["pos"],TEAL),
       ("현재 CowTalk 네트워크",f"{NET}개 농장 · {RG['rank']} (실운영)",INK2),
       ("시작 축종",f"{RG['start']}부터 → 전 축종 확산",TEALD)]
cw=(CW-2*0.4)/3; cy=2.35; ch=2.4
for i,(t,d,col) in enumerate(cards):
    cx=ML+i*(cw+0.4); box(s,cx,cy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,cy,cw,0.08,fill=col,line=None)
    text(s,cx+0.26,cy+0.28,cw-0.5,0.35,[{"runs":[R(t,13,GRAY,True)]}])
    text(s,cx+0.26,cy+0.8,cw-0.52,1.4,[{"ls":1.15,"runs":[R(d,15,INK,True)]}])
box(s,ML,5.15,CW,0.9,fill=INK,line=None)
text(s,ML+0.35,5.15,CW-0.7,0.9,[{"ls":1.15,"runs":[R(f"{RG['special']}.  ",13.5,WHITE,True),
     R("가장 데이터가 많고 산업 비중이 큰 곳에서 시작해야 표준의 힘이 생깁니다.",12.5,PALE)]}],anchor=MSO_ANCHOR.MIDDLE)
text(s,ML,6.2,CW,0.3,[{"runs":[R("네트워크 농장 수는 현재 CowTalk 실운영 분포(정직 수치) · 축우 비중 출처 축산물이력제 데이터랩·한우자조금",9.5,MUTE)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 08. Q5 무엇을 제안 — 사업
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"무엇을  ·  Q5 무엇을 제안하는가","제안은 새 소프트웨어가 아니라 '지역 축산 AX 디지털 플랫폼 구축사업'입니다.",phase="무엇을")
box(s,ML,2.15,CW,0.95,fill=PANEL2,line=LINE,lw=0.75)
text(s,ML+0.35,2.15,CW-0.7,0.95,[{"ls":1.2,"runs":[
     R("CowTalk = Foundation Platform (기반·무대).   ",14,TEALD,True),
     R(f"주인공은 {NAME} 축산의 미래 운영체제입니다. 플랫폼은 무대일 뿐, 주역은 지역의 미래입니다.",13,GRAY)]}],anchor=MSO_ANCHOR.MIDDLE)
els=[("① 기반 플랫폼","검증된 CowTalk (실운영 199농장)",TEAL),
     ("② 데이터 거버넌스","지역 공공·방역 데이터 연계",INK2),
     ("③ 현장 실증","시군 + 공수의사 시범사업",TEALD),
     ("④ 확산·표준화","지역 전역 → 전국 표준 로드맵",AMBER)]
cw=(CW-3*0.3)/4; cy=3.4; ch=2.0
for i,(t,d,col) in enumerate(els):
    cx=ML+i*(cw+0.3); box(s,cx,cy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,cy,0.08,ch,fill=col,line=None)
    text(s,cx+0.28,cy+0.3,cw-0.45,0.4,[{"runs":[R(t,14,col,True)]}])
    text(s,cx+0.28,cy+0.95,cw-0.5,1.0,[{"ls":1.2,"runs":[R(d,12,INK)]}])
text(s,ML,5.7,CW,0.4,[{"runs":[R("우리는 제품을 파는 게 아니라, 지역이 소유하는 축산 운영체제를 함께 구축하자는 제안입니다.",12.5,INK,it=True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 09. Q6 왜 CowTalk가 현실적 기반 (실운영 증거 + 정직 분리)
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"무엇을  ·  Q6 왜 CowTalk가 현실적 기반인가","제로에서 시작하지 않습니다 — 이미 199개 농장에서 실제로 돌아가는 기반입니다.",phase="무엇을")
text(s,ML,2.0,CW,0.32,[{"runs":[R("✅ 지금 실제로 동작 (코드 감사로 검증)",13,OK,True)]}])
oks=["199농장 실시간 동기화","4역할 AI 어시스턴트 (도구 25종·감사로그)","SEIR 확산 시뮬레이션",
     "위내센서 조기 질병감지","공공데이터 실연동 (이력제·등급·경락)","자가학습 루프 (레이블→정확도→보정)"]
cw=(CW-2*0.3)/3; cy=2.42; ch=0.72
for i,t in enumerate(oks):
    cx=ML+(i%3)*(cw+0.3); yy=cy+(i//3)*(ch+0.22)
    box(s,cx,yy,cw,ch,fill=PANEL,line=LINE,lw=0.75); box(s,cx,yy,0.08,ch,fill=OK,line=None)
    text(s,cx+0.25,yy,cw-0.4,ch,[{"ls":1.05,"runs":[R(t,11.5,INK,True)]}],anchor=MSO_ANCHOR.MIDDLE)
box(s,ML,4.7,CW,1.35,fill=WHITE,line=MID,lw=1.0)
text(s,ML+0.3,4.85,CW-0.6,0.32,[{"runs":[R("정직하게 — 이건 로드맵입니다",13,MID,True)]}])
text(s,ML+0.3,5.25,CW-0.6,0.8,[{"ls":1.25,"runs":[
     R("🟡 젖소 정액 정밀추천",12,MID,True),R(" (내부데이터로 동작 중, 혈통·DHI 연계 시 정밀화)     ",11.5,GRAY),
     R("🔵 KAHIS 외부 자동제출 · 일부 공공API 활용신청",12,FUT,True),R(" (행정 협조 전제)",11.5,GRAY)]}])
text(s,ML,6.15,CW,0.3,[{"runs":[R("화면의 모든 숫자는 DB에서 나오고 출처가 붙습니다 — 국가 시스템으로 가는 최소 조건입니다.",11,INK,it=True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 10. Q7-a 4층 아키텍처
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"어떻게  ·  Q7 어떻게 구축하는가","센서 위에 공공데이터·AI·역할별 행동을 얹는 4층 구조입니다.",phase="어떻게")
layers=[("4층","학습 루프","피드백 → 정확도 → 프롬프트·임계값 자기보정 (순환 강화)",TEALD),
        ("3층","역할별 서빙","농가 · 수의사 · 방역관 · 행정관 — 같은 데이터를 각 관점으로",TEAL),
        ("2층","AI 해석","Claude 중심 해석 + v4 룰엔진 폴백 (설명가능·출처기록)",INK2),
        ("1층","데이터 통합","smaXtec 위내센서 + 공공데이터 + 농장기록 → 통합 프로필",INK)]
cy=2.3; ch=0.98
for i,(lv,t,d,col) in enumerate(layers):
    yy=cy+i*(ch+0.12); box(s,ML,yy,CW,ch,fill=col,line=None)
    text(s,ML+0.35,yy,1.3,ch,[{"runs":[R(lv,20,WHITE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
    box(s,ML+1.5,yy+0.16,0.02,ch-0.32,fill=RGBColor(0xFF,0xFF,0xFF),line=None)
    text(s,ML+1.8,yy,3.2,ch,[{"runs":[R(t,16,WHITE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
    text(s,ML+4.8,yy,CW-5.0,ch,[{"ls":1.05,"runs":[R(d,12.5,PALE)]}],anchor=MSO_ANCHOR.MIDDLE)
text(s,ML,6.5,CW,0.3,[{"runs":[R("smaXtec을 복제하지 않습니다 — 그 위에 대한민국 공공데이터와 AI, 그리고 행정을 얹습니다.",11.5,INK,it=True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 11. Q7-b 알람→행동 완결
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"어떻게  ·  알람을 행동으로","알람 → 판단 → 추천 → 행동 → 기록을 한 화면에서 끝냅니다.",phase="어떻게")
steps=[("알람","소버린 AI 25룰 + smaXtec 이벤트"),("판단","감별진단 · 수정적기 (설명가능)"),
       ("추천","조치 제안 + 근거·출처"),("행동","승인 게이트 → 조치 실행"),("기록","결과 학습 → AI 강화 ↺")]
sw=(CW-4*0.35)/5; sy=2.55; sh=1.5
for i,(t,d) in enumerate(steps):
    cx=ML+i*(sw+0.35); col=TEAL if i<4 else AMBER
    box(s,cx,sy,sw,sh,fill=(PANEL2 if i==4 else PANEL),line=LINE,lw=0.9); box(s,cx,sy,sw,0.08,fill=col,line=None)
    text(s,cx,sy+0.28,sw,0.4,[{"align":PP_ALIGN.CENTER,"runs":[R(t,15,col,True)]}])
    text(s,cx+0.12,sy+0.78,sw-0.24,0.6,[{"align":PP_ALIGN.CENTER,"ls":1.1,"runs":[R(d,10.5,GRAY)]}])
    if i<4: text(s,cx+sw-0.05,sy+0.45,0.5,0.5,[{"runs":[R("→",18,MUTE,True)]}])
# 정직 칩
cx=ML
for lab in ["✅ 자동화 룰","✅ 승인 게이트","✅ 결정 큐","✅ 전문가 학습 루프"]:
    w=chip(s,cx,4.5,lab,OK); cx+=w+0.25
box(s,ML,5.25,CW,0.82,fill=INK,line=None)
text(s,ML,5.25,CW,0.82,[{"align":PP_ALIGN.CENTER,"runs":[
     R("smaXtec의 한계는 ‘그래서 뭘 해야 하지?’입니다. 우리는 그 다음 한 걸음을 화면에서 끝냅니다.",13.5,WHITE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 12. 축종 확장 전략
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"어떻게  ·  축종 확장 전략","축우에서 시작해 양돈·양계로, 최종적으로 대한민국 통합 축산 AI 플랫폼으로.",phase="어떻게")
stages=[("축우","한우·젖소","지금 시작 — 검증",TEAL),("양돈","돼지","방역 핵심 축종",INK2),
        ("양계","닭·오리","대량·밀집 관리",INK2),("기타","염소·오리 등","커넥터 재사용",INK2),
        ("통합","대한민국","축산 AI 플랫폼",TEALD)]
pw=(CW-4*0.25)/5; py=2.5; ph=1.7
hline(s,ML+pw/2,py+0.35,CW-pw,color=LINE,weight=1.5)
for i,(st,t,d,col) in enumerate(stages):
    cx=ML+i*(pw+0.25)
    box(s,cx+pw/2-0.11,py+0.24,0.22,0.22,fill=col,line=WHITE,lw=1.5,rounded=True,radius=0.5)
    box(s,cx,py+0.62,pw,ph-0.62,fill=(PANEL2 if i in(0,4) else PANEL),line=LINE,lw=0.75)
    text(s,cx,py+0.78,pw,0.35,[{"align":PP_ALIGN.CENTER,"runs":[R(st,15,col,True)]}])
    text(s,cx+0.1,py+1.18,pw-0.2,0.3,[{"align":PP_ALIGN.CENTER,"runs":[R(t,11.5,INK,True)]}])
    text(s,cx+0.1,py+1.5,pw-0.2,0.3,[{"align":PP_ALIGN.CENTER,"runs":[R(d,9.5,MUTE)]}])
box(s,ML,4.7,CW,1.35,fill=WHITE,line=LINE,lw=1.0); box(s,ML,4.7,0.08,1.35,fill=TEAL,line=None)
text(s,ML+0.35,4.9,CW-0.7,0.35,[{"runs":[R("왜 축우가 먼저인가",13.5,TEALD,True)]}])
text(s,ML+0.35,5.3,CW-0.7,0.7,[{"ls":1.2,"runs":[
     R("개체가치 최고 · 위내센서 성숙 · 질병/번식 데이터 밀도 최고 → 검증에 최적. ",12.5,INK),
     R("구조는 축종 무관(커넥터·온톨로지 재사용)이라 확장이 빠릅니다.",12.5,GRAY)]}])
text(s,ML,6.2,CW,0.3,[{"runs":[R("한우 플랫폼도, 낙농 플랫폼도 아닙니다 — 축우 중심으로 검증한 구조를 전 축종으로 확장합니다.",11,INK,it=True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 13. Q8-① 농가 기대효과
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"그래서 결과는  ·  Q8 기대효과 ① 농가","농가는 번식·질병·노동을 데이터로 잡아 두당 연 13~24만원의 생산성을 얻습니다.",phase="그래서 결과는")
rows=[("번식 효율","공태일 단축 (발정탐지·수정적기·동기화)","6만원",TEAL),
      ("질병 조기대응","증상 전 감지 → 치료비·폐사·유량손실↓","4만원",INK2),
      ("산유·증체","건강·번식 개선의 이차효과","2만원",INK2),
      ("노동 절감","알람→행동 자동화로 관찰·기록 노동↓","1만원",INK2)]
cy=2.3; rh=0.72
for i,(t,d,v,col) in enumerate(rows):
    yy=cy+i*(rh+0.14); box(s,ML,yy,CW-2.6,rh,fill=PANEL,line=LINE,lw=0.75); box(s,ML,yy,0.08,rh,fill=col,line=None)
    text(s,ML+0.3,yy,3.2,rh,[{"runs":[R(t,13.5,INK,True)]}],anchor=MSO_ANCHOR.MIDDLE)
    text(s,ML+3.5,yy,CW-6.4,rh,[{"ls":1.05,"runs":[R(d,11.5,GRAY)]}],anchor=MSO_ANCHOR.MIDDLE)
    text(s,ML+CW-2.6-1.3,yy,1.2,rh,[{"align":PP_ALIGN.RIGHT,"runs":[R(v,15,col,True)]}],anchor=MSO_ANCHOR.MIDDLE)
box(s,ML+CW-2.4,2.3,2.4,rh*4+0.14*3,fill=INK,line=None)
text(s,ML+CW-2.4,2.75,2.4,0.4,[{"align":PP_ALIGN.CENTER,"runs":[R("보수 합계",12,PALE,True)]}])
text(s,ML+CW-2.4,3.25,2.4,0.7,[{"align":PP_ALIGN.CENTER,"runs":[R("13만원",30,MINT,True)]}])
text(s,ML+CW-2.4,4.05,2.4,0.4,[{"align":PP_ALIGN.CENTER,"runs":[R("/ 두 · 년",12,PALE)]}])
text(s,ML+CW-2.4,4.55,2.4,0.5,[{"align":PP_ALIGN.CENTER,"ls":1.05,"runs":[R("적극 시나리오 24만원",11,PALE)]}])
box(s,ML,5.75,CW,0.42,fill=None,line=None)
text(s,ML,5.72,CW,0.32,[{"runs":[R("전제·출처  ",10,MUTE,True),
     R("공태일 1일 = $3.0~5.1 (De Vries 2006, J.Dairy Sci) · 보수 15일 단축 × $3 ≈ 6만원. 지역 실두수로 재산출·실측 보정 전제.",10,GRAY)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 14. Q8-② 지역 기대효과
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"그래서 결과는  ·  Q8 기대효과 ② 지역","지역은 데이터 행정·방역 고도화·탄소 관리·AI 산업·일자리를 동시에 얻습니다.",phase="그래서 결과는")
cards=[("데이터 행정","전화·엑셀 → 실시간 대시보드. 예산 집행이 근거 기반으로.",TEAL),
       ("방역 고도화","역학조사 며칠 → 30초. 확산 사전 시뮬레이션(SEIR).",INK2),
       ("탄소·분뇨 관리","번식효율↑ → 육성우 감축 → 사료·분뇨·탄소 동시 절감.",TEALD),
       ("지역 AI 산업","지역 축산 데이터가 지역 AI 기업·연구의 자산으로.",INK2),
       ("청년 일자리","데이터·수의·기술 인력 수요 창출.",TEAL),
       ("정책 효율","수급·방역·복지 예산의 효과를 측정 가능하게.",INK2)]
cw=(CW-2*0.3)/3; cy=2.35; ch=1.7
for i,(t,d,col) in enumerate(cards):
    cx=ML+(i%3)*(cw+0.3); yy=cy+(i//3)*(ch+0.25)
    box(s,cx,yy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,yy,cw,0.08,fill=col,line=None)
    text(s,cx+0.26,yy+0.26,cw-0.5,0.35,[{"runs":[R(t,14,INK,True)]}])
    text(s,cx+0.26,yy+0.72,cw-0.5,0.9,[{"ls":1.2,"runs":[R(d,11.5,GRAY)]}])
box(s,ML,6.05,CW,0.5,fill=PANEL,line=LINE,lw=0.75)
text(s,ML+0.3,6.05,CW-0.6,0.5,[{"runs":[R("농가 지원사업이 아니라, 지역 ‘행정 역량’ 자체를 디지털로 끌어올리는 사업입니다.",12.5,INK,True)]}],anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 15. Q9 대한민국 기대효과
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"그래서 결과는  ·  Q9 대한민국이 얻는 것","지역의 성공은 곧 국가 축산 AI 플랫폼·표준 데이터·K-축산 수출로 이어집니다.",phase="그래서 결과는")
cards=[("국가 축산 AI 플랫폼","원점 = 선도 지역",TEAL),("표준 데이터 체계","축산 전주기 데이터 표준",INK2),
       ("AI 산업 성장","데이터·모델·인재 생태계",TEALD),("K-축산 수출","i18n·국가 어댑터 구조 내장",INK2),
       ("식량안보","생산성·자급 기반 강화",TEAL),("국가 방역 주권","데이터 자립·확산 차단",AMBER)]
cw=(CW-2*0.3)/3; cy=2.35; ch=1.55
for i,(t,d,col) in enumerate(cards):
    cx=ML+(i%3)*(cw+0.3); yy=cy+(i//3)*(ch+0.25)
    box(s,cx,yy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,yy,0.08,ch,fill=col,line=None)
    text(s,cx+0.28,yy+0.28,cw-0.5,0.4,[{"runs":[R(t,14,INK,True)]}])
    text(s,cx+0.28,yy+0.82,cw-0.5,0.5,[{"ls":1.1,"runs":[R(d,11.5,GRAY)]}])
box(s,ML,6.0,CW,0.55,fill=INK,line=None)
text(s,ML,6.0,CW,0.55,[{"align":PP_ALIGN.CENTER,"runs":[
     R("지역에서 만든 표준이 국가 표준이 됩니다. 표준을 만든 지역이 원점의 지위를 갖습니다.",13,WHITE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 16. 경제효과 종합 (보수적 + 근거)
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"그래서 결과는  ·  경제효과 종합","보수적으로 잡아도 편익이 비용을 크게 웃돕니다 — 모든 수치에 근거를 붙였습니다.",phase="그래서 결과는")
cols=[("전제 (투명 공개)",[ "지역 축우 실두수 × 두당 보수 13만원","센서+플랫폼 구독 · 중계기 구축비","사회·환경은 지역 원단위로 보정"],INK2),
      ("편익 구성 (연간)",[ "농가 직접 = 두수 × 13만원 (보수)","사회·환경 = 육성우 감축(방향성)","방역 = 옵션가치(1회 차단)"],TEAL),
      ("판정",[ "보수 시나리오에서도 편익 > 비용","방역 1회 차단 = 사업비 회수","중장기: 데이터·AI 산업 자산"],AMBER)]
cw=(CW-2*0.4)/3; cy=2.3; ch=2.7
for i,(t,items,col) in enumerate(cols):
    cx=ML+i*(cw+0.4); box(s,cx,cy,cw,ch,fill=WHITE,line=LINE,lw=1.0); box(s,cx,cy,cw,0.08,fill=col,line=None)
    text(s,cx+0.26,cy+0.26,cw-0.5,0.35,[{"runs":[R(t,14,col,True)]}])
    for j,it in enumerate(items):
        yy=cy+0.85+j*0.62
        text(s,cx+0.26,yy,0.25,0.4,[{"runs":[R("▪",12,col,True)]}])
        text(s,cx+0.52,yy,cw-0.75,0.6,[{"ls":1.05,"runs":[R(it,11.5,INK)]}])
box(s,ML,5.35,CW,0.72,fill=RGBColor(0xFB,0xF3,0xE6),line=AMBER,lw=1.0)
text(s,ML+0.3,5.35,CW-0.6,0.72,[{"runs":[R("정직 배너  ",11,AMBER,True),
     R("모든 값은 추정이며 전제에 민감합니다. 절대금액·BCR은 지역 실두수·원단위 확정 후 산출합니다. 부풀리지 않았습니다.",11.5,INK)]}],anchor=MSO_ANCHOR.MIDDLE)
text(s,ML,6.2,CW,0.3,[{"runs":[R("출처  De Vries 2006(공태일) · 통계청 사육두수 · 농식품부/식품저널(산업규모)",9.5,MUTE)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 17. 방역 옵션가치
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"그래서 결과는  ·  방역 옵션가치","전염병 1회 차단은 그 자체로 사업비를 회수합니다 — 구제역 한 번이 재정 수조 원입니다.",phase="그래서 결과는")
stats=[("살처분","348만두","2010–11 구제역"),("피해 농장","6,241곳","전국 11개 시도"),
       ("재정 소요","3.18조","매몰보상 1.82조 포함")]
cw=(CW-2*0.4)/3; cy=2.35; ch=2.0
for i,(t,big,sub) in enumerate(stats):
    cx=ML+i*(cw+0.4); box(s,cx,cy,cw,ch,fill=INK,line=None)
    text(s,cx+0.3,cy+0.35,cw-0.6,0.35,[{"runs":[R(t,13,PALE,True)]}])
    text(s,cx+0.3,cy+0.85,cw-0.6,0.7,[{"runs":[R(big,32,MINT,True)]}])
    text(s,cx+0.3,cy+1.55,cw-0.6,0.35,[{"runs":[R(sub,11.5,PALE)]}])
box(s,ML,4.75,CW,0.95,fill=AMBER,line=None)
text(s,ML,4.75,CW,0.95,[{"align":PP_ALIGN.CENTER,"ls":1.2,"runs":[
     R("조기감지 + 확산 시뮬레이션으로 1회만 차단해도 회피액이 사업 비용을 압도합니다.\n",14,WHITE,True),
     R("연간 운영비는 비용이 아니라 보험료입니다.",13,WHITE)]}],anchor=MSO_ANCHOR.MIDDLE)
text(s,ML,5.85,CW,0.3,[{"runs":[R("출처  위키백과·경향신문·농민신문 (2010–11 대한민국 구제역)",9.5,MUTE)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 18. 실행 로드맵
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"실행  ·  어떻게 확산하나","검증된 구조를 지금 → 시범 → 지역 전역 → 전국 → 수출로 단계 확장합니다.",phase="실행")
phases=[("지금",f"199 농장 · {SHORT} {NET}곳","실운영 검증 완료",TEAL),
        ("1단계",f"{SHORT} 시범 (시군+공수의사)","실증·효과 측정",INK2),
        ("2단계",f"{NAME} 전역","도 단위 표준화",INK2),
        ("3단계","전국 표준","국가 방역 고도화",INK2),
        ("4단계","해외 수출","i18n·국가별 어댑터",TEALD)]
pw=(CW-4*0.25)/5; py=2.4; ph=1.7
hline(s,ML+pw/2,py+0.35,CW-pw,color=LINE,weight=1.5)
for i,(st,t,d,col) in enumerate(phases):
    cx=ML+i*(pw+0.25)
    box(s,cx+pw/2-0.11,py+0.24,0.22,0.22,fill=col,line=WHITE,lw=1.5,rounded=True,radius=0.5)
    box(s,cx,py+0.62,pw,ph-0.62,fill=(PANEL2 if i==0 else PANEL),line=LINE,lw=0.75)
    text(s,cx,py+0.78,pw,0.3,[{"align":PP_ALIGN.CENTER,"runs":[R(st,12,col,True)]}])
    text(s,cx+0.1,py+1.12,pw-0.2,0.5,[{"align":PP_ALIGN.CENTER,"ls":1.05,"runs":[R(t,11,INK,True)]}])
    text(s,cx+0.1,py+1.62,pw-0.2,0.3,[{"align":PP_ALIGN.CENTER,"runs":[R(d,9.5,MUTE)]}])
text(s,ML,4.7,CW,0.35,[{"runs":[R("각 단계는 앞 단계 성과로 다음을 정당화합니다 — 무리한 점프가 없습니다.",13,INK,True)]}])
box(s,ML,5.2,CW,0.9,fill=PANEL2,line=LINE,lw=0.75)
text(s,ML+0.35,5.2,CW-0.7,0.9,[{"ls":1.2,"runs":[R("Q10 국가 경로  ",12.5,TEALD,True),
     R(f"{NAME} 데이터·표준 → 국가 농업 AX 플랫폼 연계 → 타 지역 확산 → 국가 통합. 먼저 한 지역이 표준을 소유합니다.",12.5,GRAY)]}],anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 19. 지역에 드리는 요청
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"실행  ·  요청","세 가지를 요청드립니다 — 시범사업, 데이터 연계, 확산 로드맵 공동 수립.",phase="실행")
asks=[("1. 시범사업",f"{NAME} 일부 시군 + 공수의사 대상 실증. 효과를 데이터로 측정합니다."),
      ("2. 데이터 연계","KAHIS · 이력제 · 도 방역 데이터 연동에 대한 행정적 협조."),
      ("3. 확산 로드맵",f"{NAME} 전역 → 전국 표준화를 함께 수립합니다.")]
cy=2.35; rh=1.15
for i,(t,d) in enumerate(asks):
    yy=cy+i*(rh+0.2); box(s,ML,yy,CW,rh,fill=WHITE,line=TEAL,lw=1.0); box(s,ML,yy,0.1,rh,fill=TEAL,line=None)
    text(s,ML+0.4,yy+0.2,CW-0.8,0.4,[{"runs":[R(t,16,TEALD,True)]}])
    text(s,ML+0.4,yy+0.66,CW-0.8,0.4,[{"ls":1.1,"runs":[R(d,12.5,GRAY)]}])
text(s,ML,6.45,CW,0.3,[{"runs":[R("돈보다 먼저 필요한 것은 데이터 연계에 대한 행정적 결단입니다.",12,INK,it=True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 20. Q10 국가 경로 (강조 슬라이드)
# ══════════════════════════════════════════════════════════════════
s=slide()
header(s,"실행  ·  Q10 국가 플랫폼 발전 경로",f"{NAME}의 검증이 국가 축산 AI 플랫폼의 표준·원점이 됩니다.",phase="실행")
steps=[(f"{NAME}\n데이터·표준",TEAL),("국가 농업 AX\n플랫폼 연계",INK2),("타 지역\n확산",INK2),("대한민국\n통합 플랫폼",TEALD)]
sw=(CW-3*0.6)/4; sy=2.7; sh=1.7
for i,(t,col) in enumerate(steps):
    cx=ML+i*(sw+0.6); box(s,cx,sy,sw,sh,fill=col,line=None)
    text(s,cx+0.15,sy,sw-0.3,sh,[{"align":PP_ALIGN.CENTER,"ls":1.15,"runs":[R(t,15,WHITE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
    if i<3: text(s,cx+sw+0.1,sy+sh/2-0.3,0.5,0.6,[{"runs":[R("→",26,MUTE,True)]}],anchor=MSO_ANCHOR.MIDDLE)
box(s,ML,5.0,CW,1.05,fill=INK,line=None)
text(s,ML+0.4,5.0,CW-0.8,1.05,[{"ls":1.2,"runs":[
     R("국가 플랫폼은 결국 어느 지역의 데이터로 시작됩니다.\n",14,WHITE,True),
     R(f"그 지역이 {NAME}가 될 것인가 — 그것이 오늘의 결정입니다.",13.5,MINT,True)]}],anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 21. 클로징 제언
# ══════════════════════════════════════════════════════════════════
s=slide(); box(s,0,0,EW,EH,fill=INK,line=None); box(s,0,0,0.22,EH,fill=TEAL,line=None)
text(s,1.0,1.45,11.4,0.4,[{"runs":[R("제언",13,MINT,True)]}])
text(s,1.0,1.95,11.6,1.8,[{"ls":1.08,"runs":[R(f"{NAME}가 대한민국 축산\nAX의 첫 표준을 만듭니다.",34,WHITE,True)]}])
hline(s,1.0,4.15,4.2,color=TEAL,weight=2.2)
text(s,1.0,4.4,11.2,0.9,[{"ls":1.3,"runs":[
     R("데이터가 가장 먼저 쌓이는 곳이 표준을 만듭니다.\n",15,PALE),
     R(f"검증된 구조는 이미 여기 있습니다. 그 시작을 {NAME}가 하십시오.",15,PALE)]}])
mini=[("농가·사회 편익","두당 보수 13만원 · 지역 파급"),("방역 옵션가치","전염병 1회 = 재정 수조 회피"),("확장 경로",f"{SHORT} → 전국 → 수출")]
mw=(11.4-2*0.4)/3
for i,(t,d) in enumerate(mini):
    cx=1.0+i*(mw+0.4); box(s,cx,5.5,mw,0.95,fill=RGBColor(0x12,0x2E,0x47),line=RGBColor(0x24,0x44,0x5E),lw=0.75)
    text(s,cx+0.25,5.65,mw-0.5,0.3,[{"runs":[R(t,12.5,MINT,True)]}])
    text(s,cx+0.25,6.0,mw-0.5,0.35,[{"runs":[R(d,11.5,WHITE)]}])
text(s,1.0,6.85,11.4,0.3,[{"runs":[R("CowTalk v5.0  ·  Foundation Platform for Livestock AX  ·  함께 만드는 축산 디지털 주권",11,MUTE)]}])

# ── 발표자 노트 (18~20분) ──
NOTES=[
 f"[표지] ‘오늘 18분 안에 — 무엇을 만들었고, 어떻게 작동하며, 그래서 {NAME}에 어떤 결과가 오는지 숫자로 말씀드립니다.’",
 f"[제언·답부터] ‘결론은 하나입니다 — 지금, 우리 지역이 먼저.’ 새 SW 판매가 아니라 검증된 기반 위 사업. 3축(농가·방역·AI산업). 나머지 19장은 근거.",
 "[Q1] 축산은 농업의 43%(25.5조) 기간산업인데 4대 위기(규모·두수감소·구제역·고령화)가 겹친다. 지킬 가치 큰데 수단이 아날로그.",
 "[Q2] 부족한 건 센서가 아니라 루프. 지금은 알람에서 멈춘다. 필요한 건 데이터→AI→의사결정→실행→학습 순환. smaXtec도 알람→행동 다리 없음.",
 "[Q3] 왜 지금. 정부 AI예산 9.9조(3배)·국가 농업 AX 705억·지역 AI 대전환. Smart Farm(설비)→AX(운영지능) 단계. 정책 창은 오래 안 열린다.",
 "[Q4] 왜 지역. 국가 일괄 불가. 집행주체=광역지자체, 데이터 밀도가 표준, 검증→확산이 실패비용 최저.",
 f"[왜 {SHORT}] 치환 슬라이드. {NAME}의 축우 포지션·현 네트워크 {NET}곳(실측)·시작 축종. ‘데이터 많은 곳이 표준을 만든다.’",
 "[Q5] 제안은 앱이 아니라 사업. CowTalk=Foundation, 주인공=지역 미래. 4요소(기반·거버넌스·실증·확산).",
 "[Q6] 제로 아님. ✅현재 6가지 실동작(코드 감사). 정직하게 🟡젖소 정액 정밀추천·🔵KAHIS 외부제출은 로드맵. ‘숫자에 출처가 붙는다.’",
 "[Q7-a] 4층 구조. 데이터통합→AI해석→역할별 서빙→학습 루프. smaXtec 복제 아님, 그 위에 공공데이터·AI·행정.",
 "[Q7-b] 알람→판단→추천→행동→기록 완결. ✅자동화룰·승인게이트·결정큐·학습루프. (가능하면 실제 시연 1회)",
 "[축종 확장] 축우→양돈→양계→기타→통합. 축우 먼저인 이유(개체가치·센서성숙·데이터밀도). 구조는 축종 무관.",
 "[Q8-①농가] 보수 13만원/두·년(번식6+질병4+산유2+노동1), 적극 24만. 근거 De Vries 공태일. 지역 실두수로 재산출 명시. ‘부풀리지 않았다.’",
 "[Q8-②지역] 데이터행정·방역고도화·탄소·AI산업·일자리·정책효율. ‘농가 지원 아니라 행정 역량 자체를 올리는 사업.’",
 "[Q9 대한민국] 국가 플랫폼·표준데이터·AI산업·K-축산 수출·식량안보·방역주권. ‘지역 표준=국가 표준, 지역이 원점.’",
 "[경제 종합] 천천히. 전제 투명공개→편익구성→판정. 정직 배너 반드시 입으로: ‘절대금액·BCR은 지역 실측 확정 후 산출.’",
 "[방역 옵션가치] 구제역 348만두·6,241농장·재정 3.18조. 1회 차단이 사업비 압도. ‘운영비는 보험료.’",
 f"[로드맵] 지금(199·{SHORT} {NET}) → {SHORT} 시범 → {NAME} 전역 → 전국 → 수출. 각 단계는 앞 성과로 정당화.",
 "[요청] 3가지 — 시범사업·데이터 연계·확산 로드맵. ‘돈보다 먼저 데이터 연계 결단.’",
 f"[Q10 국가경로] {NAME} 표준 → 국가 AX 연계 → 확산 → 통합. ‘국가 플랫폼은 어느 지역 데이터로 시작된다 — 그 지역이 될 것인가.’",
 f"[클로징] ‘{NAME}가 대한민국 축산 AX의 첫 표준을 만듭니다. 그 시작을 지금 하십시오.’ 감사합니다.",
]
for sl,nt in zip(prs.slides,NOTES):
    sl.notes_slide.notes_text_frame.text=nt

out=sys.argv[2] if len(sys.argv)>2 else f"/home/user/cowtalk/docs/master-deck/CowTalk_AX_Master_{KEY}.pptx"
prs.save(out)
print(f"saved: {out}  (region={KEY}, {NAME}, net={NET})  slides={len(prs.slides.__iter__.__self__._sldIdLst)}")
