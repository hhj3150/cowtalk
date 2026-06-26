# -*- coding: utf-8 -*-
# CowTalk × 경기도 인수위 — 15분 브리핑 덱 (McKinsey 스타일, simple is best)
# 논리 골격: 무엇을 → 어떻게 → 그래서 결과는 / 액션 타이틀(결론부터)
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 디자인 토큰 ──
INK   = RGBColor(0x0B,0x22,0x39)   # 네이비 잉크
INK2  = RGBColor(0x24,0x3B,0x53)
TEAL  = RGBColor(0x0E,0x7C,0x7B)   # 액센트
TEALD = RGBColor(0x09,0x53,0x52)
MINT  = RGBColor(0x16,0xB0,0xA6)
AMBER = RGBColor(0xC9,0x7A,0x1E)   # 결과/금액 강조
GRAY  = RGBColor(0x55,0x66,0x76)   # 보조 텍스트
MUTE  = RGBColor(0x90,0x9E,0xAA)
LINE  = RGBColor(0xD8,0xE0,0xE6)   # 헤어라인
PANEL = RGBColor(0xF4,0xF7,0xF8)   # 연한 패널
PANEL2= RGBColor(0xEA,0xF1,0xF1)   # 연한 틸 패널
WHITE = RGBColor(0xFF,0xFF,0xFF)
KR = "맑은 고딕"

prs = Presentation()
prs.slide_width  = I(13.333)
prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
EW, EH = 13.333, 7.5
ML, MR = 0.72, 0.72
CW = EW - ML - MR

def slide():
    return prs.slides.add_slide(BLANK)

def _set_fill(shape, color):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = color

def _set_line(shape, color, w=0.75):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color; shape.line.width = Pt(w)

def box(s, x,y,w,h, fill=None, line=None, lw=0.75, rounded=False, radius=0.06):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        I(x),I(y),I(w),I(h))
    _set_fill(shp, fill); _set_line(shp, line, lw)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp

def hline(s, x,y,w, color=LINE, weight=0.75):
    ln = s.shapes.add_connector(2, I(x),I(y),I(x+w),I(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def text(s, x,y,w,h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: list of paragraph dicts {align, sb, sa, ls, runs:[{t,sz,c,b,f,it}]}"""
    tb = s.shapes.add_textbox(I(x),I(y),I(w),I(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, 0)
    if isinstance(paras, str):
        paras = [{"runs":[{"t":paras}]}]
    for i,pa in enumerate(paras):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        if pa.get("sb") is not None: p.space_before = Pt(pa["sb"])
        if pa.get("sa") is not None: p.space_after  = Pt(pa["sa"])
        if pa.get("ls") is not None: p.line_spacing = pa["ls"]
        for r in pa["runs"]:
            run = p.add_run(); run.text = r["t"]
            f = run.font
            f.name = r.get("f", KR); f.size = Pt(r.get("sz",14))
            f.bold = r.get("b", False); f.italic = r.get("it", False)
            f.color.rgb = r.get("c", INK)
            # CJK 폰트 강제 (eastAsia)
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
            ea.set('typeface', r.get("f", KR))
    return tb

def R(t, sz=14, c=INK, b=False, f=KR, it=False):
    return {"t":t,"sz":sz,"c":c,"b":b,"f":f,"it":it}

PHASE_STYLE = {
    "현황":      (RGBColor(0x6B,0x7A,0x88), WHITE),
    "무엇을":    (TEAL, WHITE),
    "어떻게":    (INK, WHITE),
    "그래서 결과는": (AMBER, WHITE),
    "실행":      (INK, WHITE),
    "제언":      (TEALD, WHITE),
}

def header(s, kicker, title, phase=None, title_sz=23):
    text(s, ML, 0.46, CW-2.6, 0.3,
         [{"runs":[R(kicker, 11.5, TEAL, True)]}])
    if phase:
        bg, fg = PHASE_STYLE.get(phase,(TEAL,WHITE))
        pw = 1.05 + 0.115*len(phase)
        pb = box(s, EW-MR-pw, 0.44, pw, 0.34, fill=bg, line=None, rounded=True, radius=0.5)
        text(s, EW-MR-pw, 0.44, pw, 0.34,
             [{"align":PP_ALIGN.CENTER,"runs":[R(phase, 11, fg, True)]}],
             anchor=MSO_ANCHOR.MIDDLE)
    text(s, ML, 0.80, CW, 1.0,
         [{"ls":1.04,"runs":[R(title, title_sz, INK, True)]}])
    hline(s, ML, 1.86, 0.62, color=TEAL, weight=2.6)

TOTAL = 14
_page = [1]   # 표지=01(footer 없음). footer() 호출마다 자동 증가.
def footer(s):
    _page[0] += 1; n = _page[0]
    hline(s, ML, 7.04, CW, color=LINE, weight=0.75)
    text(s, ML, 7.10, 8.0, 0.3,
         [{"runs":[R("CowTalk", 9, TEAL, True), R("  ×  경기도 · 축산 디지털 행정 대전환", 9, MUTE)]}])
    text(s, EW-MR-2.0, 7.10, 2.0, 0.3,
         [{"align":PP_ALIGN.RIGHT,"runs":[R(f"{n:02d} / {TOTAL}", 9, MUTE, True)]}])

# ══════════════════════════════════════════════════════════════════
# 1. 표지
# ══════════════════════════════════════════════════════════════════
s = slide()
box(s, 0,0, EW, EH, fill=INK, line=None)
box(s, 0,0, 0.22, EH, fill=TEAL, line=None)
# 우상단 미세 그래픽 (레이어 라인)
for i,yy in enumerate([1.0,1.34,1.68]):
    hline(s, EW-3.4, yy, 2.7, color=RGBColor(0x1E,0x3A,0x52), weight=1.4)
text(s, 1.0, 1.6, 11, 0.4, [{"runs":[R("경기도 인수위원회 정책 브리핑 · 2026", 13, MINT, True)]}])
text(s, 1.0, 2.25, 11.4, 2.2,
     [{"ls":1.02,"runs":[R("축산 디지털 행정 대전환", 46, WHITE, True)]},
      {"sb":10,"ls":1.1,"runs":[R("센서 · 공공데이터 · AI로 완성하는 경기도 축산 운영체제", 18, RGBColor(0xC7,0xD6,0xDE))]}])
hline(s, 1.0, 4.95, 4.2, color=TEAL, weight=2.2)
text(s, 1.0, 5.15, 11.4, 0.8,
     [{"runs":[R("“경기도가 대한민국 축산 디지털 행정을 선도한다.”", 17, WHITE, True, it=True)]}])
text(s, 1.0, 6.7, 11.4, 0.4,
     [{"runs":[R("CowTalk v5.0", 12, WHITE, True), R("   ·   D2O Corp × smaXtec Korea", 12, MUTE)]}])

# ══════════════════════════════════════════════════════════════════
# 2. 제언 (답부터)
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "제언  ·  THE ANSWER", "경기도가 대한민국 축산 디지털 행정 전환을 선도해야 합니다.", phase="제언", title_sz=24)
pillars = [
 ("01", "방역 주권", TEAL,
  "전염병 1회 차단 = 수십~수백억.\n역학조사 며칠 → 30초.\n전국 통합 실시간 감시.", "옵션가치 회수형"),
 ("02", "농가 소득", INK2,
  "번식·질병을 데이터로 관리.\n젖소 두당 생산성 +24만원/년.\n경기 농가 편익 240억/년.", "+240억 / 년"),
 ("03", "사회·환경", TEALD,
  "육성우 감축 — 적게 키우고\n더 잘 짠다. 사료·분뇨·탄소\n동시 절감 626억/년.", "+626억 / 년"),
]
cardw, gap = (CW-2*0.4)/3, 0.4
cy = 2.18; ch = 3.05
for i,(num,tt,col,body,tag) in enumerate(pillars):
    cx = ML + i*(cardw+gap)
    box(s, cx, cy, cardw, ch, fill=WHITE, line=LINE, lw=1.0)
    box(s, cx, cy, cardw, 0.08, fill=col, line=None)
    text(s, cx+0.28, cy+0.30, cardw-0.5, 0.6, [{"runs":[R(num, 26, RGBColor(0xCD,0xD7,0xDD), True)]}])
    text(s, cx+0.28, cy+0.92, cardw-0.5, 0.4, [{"runs":[R(tt, 18, INK, True)]}])
    text(s, cx+0.28, cy+1.45, cardw-0.5, 1.2, [{"ls":1.18,"runs":[R(body, 12.5, GRAY)]}])
    box(s, cx+0.28, cy+ch-0.62, cardw-0.56, 0.4, fill=PANEL2, line=None, rounded=True, radius=0.4)
    text(s, cx+0.28, cy+ch-0.62, cardw-0.56, 0.4,
         [{"align":PP_ALIGN.CENTER,"runs":[R(tag, 12, TEALD, True)]}], anchor=MSO_ANCHOR.MIDDLE)
# 하단 ask bar
by = 5.55
box(s, ML, by, CW, 0.92, fill=INK, line=None)
text(s, ML+0.45, by, 8.1, 0.92,
     [{"runs":[R("투자 ", 14, RGBColor(0xC7,0xD6,0xDE)), R("180억/년", 19, WHITE, True),
               R("  →  편익 ", 14, RGBColor(0xC7,0xD6,0xDE)), R("896억/년", 19, MINT, True),
               R("   ·   BCR ≈ 5.0", 14, RGBColor(0xC7,0xD6,0xDE))]}], anchor=MSO_ANCHOR.MIDDLE)
box(s, ML+8.8, by+0.2, 0.02, 0.52, fill=RGBColor(0x2A,0x47,0x60), line=None)
text(s, ML+9.05, by, CW-9.05-0.3, 0.92,
     [{"align":PP_ALIGN.RIGHT,"runs":[R("경기가 첫 표준을 선점", 14, MINT, True)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 2-B. 경기도 축산 현황 — 왜 경기인가 (낙농 중심)
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "현황  ·  왜 경기인가", "대한민국 낙농의 중심은 경기도 — 축우에서 시작해 전 축종으로 확장합니다.", phase="현황")
gcards=[("41%","경기 젖소 사육 점유","전국 최대 낙농 거점",TEAL),
        ("50%+","수도권 인구 집중","최대 소비지 = 경기",INK2),
        ("낙농 1번지","생산 × 소비 동시 중심","유일한 수도권 광역",TEALD)]
cardw=(CW-2*0.4)/3; cy=2.12; ch=1.9
for i,(big,tt,sub,col) in enumerate(gcards):
    cx=ML+i*(cardw+0.4)
    box(s, cx,cy,cardw,ch, fill=WHITE, line=LINE, lw=1.0)
    box(s, cx,cy,cardw,0.08, fill=col, line=None)
    text(s, cx+0.28, cy+0.34, cardw-0.5, 0.8, [{"runs":[R(big, (40 if i<2 else 25), col, True)]}])
    text(s, cx+0.28, cy+1.2, cardw-0.5, 0.35, [{"runs":[R(tt, 14, INK, True)]}])
    text(s, cx+0.28, cy+1.54, cardw-0.5, 0.3, [{"runs":[R(sub, 11, MUTE)]}])
# 전략 흐름: 축우 시작 → 타 축종 확산
text(s, ML, 4.2, CW, 0.3, [{"runs":[R("디지털 대전환 전개 — ", 13, GRAY), R("축우 중심으로 시작 → 전 축종 확산", 13, INK, True)]}])
fy=4.58; fh=0.95
box(s, ML, fy, 4.7, fh, fill=TEAL, line=None)
text(s, ML, fy, 4.7, fh, [{"align":PP_ALIGN.CENTER,"ls":1.1,"runs":[R("축우 — 젖소 · 한우\n축산 디지털 대전환 시작", 14, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, ML+4.7, fy, 0.9, fh, [{"align":PP_ALIGN.CENTER,"runs":[R("확산 ›", 16, AMBER, True)]}], anchor=MSO_ANCHOR.MIDDLE)
exp=["양돈","가금","기타 축종"]
ex0=ML+5.7; exw=(EW-MR-ex0-2*0.22)/3
for i,t in enumerate(exp):
    cx=ex0+i*(exw+0.22)
    box(s, cx, fy, exw, fh, fill=PANEL, line=LINE, lw=0.75)
    text(s, cx, fy, exw, fh, [{"align":PP_ALIGN.CENTER,"runs":[R(t, 14, GRAY, True)]}], anchor=MSO_ANCHOR.MIDDLE)
# takeaway
box(s, ML, 5.78, CW, 0.78, fill=INK, line=None)
text(s, ML, 5.78, CW, 0.78, [{"align":PP_ALIGN.CENTER,"runs":[R("사육도 소비도 경기 — 경기가 낙농을 디지털로 표준화하면, 그것이 곧 대한민국 표준이 됩니다.", 14, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 3. 현황 (Situation)
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "현황  ·  지금 어디에 서 있나", "세계는 축산을 데이터로 운영하지만, 방역·행정은 아직 전화와 엑셀입니다.", phase="현황")
cols = [
 ("농가", "센서는 도입했지만\n“알람이 와도 무엇을 할지”는\n각자 경험으로 판단."),
 ("방역", "역학조사는 전화·수기 엑셀.\n접촉 농장 파악에 며칠.\n확산은 늘 사후 대응."),
 ("행정", "통계는 아날로그·취합형.\n수급·방역 의사결정이\n데이터보다 늦게 도착."),
]
cardw = (CW-2*0.4)/3
cy=2.25; ch=2.55
for i,(tt,body) in enumerate(cols):
    cx = ML+i*(cardw+0.4)
    box(s, cx,cy,cardw,ch, fill=PANEL, line=LINE, lw=1.0)
    text(s, cx+0.28, cy+0.26, cardw-0.5, 0.4, [{"runs":[R(tt, 16, TEAL, True)]}])
    hline(s, cx+0.28, cy+0.74, 0.5, color=TEAL, weight=2.0)
    text(s, cx+0.28, cy+0.92, cardw-0.5, 1.5, [{"ls":1.22,"runs":[R(body, 13.5, INK)]}])
box(s, ML, 5.3, CW, 0.8, fill=PANEL2, line=None)
text(s, ML+0.35, 5.3, CW-0.7, 0.8,
     [{"runs":[R("→  데이터는 쌓이는데, ", 15, INK), R("판단과 행동으로 이어지지 않는다.", 15, TEALD, True)]}],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 4. 문제 (Complication)
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "현황  ·  핵심 공백", "세계 최고 센서(smaXtec)도 ‘알람 → 행동’의 다리가 없습니다.", phase="현황")
# 흐름 다이어그램
fy=2.55; fh=1.15
b1w, b3w = 3.5, 3.9; gapw = (CW-b1w-b3w)
box(s, ML, fy, b1w, fh, fill=INK, line=None)
text(s, ML, fy, b1w, fh, [{"align":PP_ALIGN.CENTER,"ls":1.1,"runs":[R("smaXtec 위내센서\n발정·발열 알람", 15, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
# 끊긴 화살표 + 물음표
text(s, ML+b1w, fy, gapw, fh, [{"align":PP_ALIGN.CENTER,"runs":[R("─ ─ ─ ✕ ─ ─ ─", 18, AMBER, True)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, ML+b1w, fy+fh-0.02, gapw, 0.4, [{"align":PP_ALIGN.CENTER,"runs":[R("연결 끊김", 11, AMBER, True)]}], anchor=MSO_ANCHOR.TOP)
box(s, EW-MR-b3w, fy, b3w, fh, fill=WHITE, line=AMBER, lw=1.4)
text(s, EW-MR-b3w, fy, b3w, fh, [{"align":PP_ALIGN.CENTER,"ls":1.1,"runs":[R("농가·방역관:\n“그래서 뭘 해야 하지?”", 15, AMBER, True)]}], anchor=MSO_ANCHOR.MIDDLE)
# 3 fact
facts = ["알람은 정확해도 다음 행동은\n사람이 혼자 판단", "추천·기록·역학으로\n이어지는 워크플로 부재", "한국 공공데이터(이력·방역)와\n연결 안 됨"]
cardw=(CW-2*0.4)/3; cy=4.35; ch=1.15
for i,b in enumerate(facts):
    cx=ML+i*(cardw+0.4)
    box(s, cx,cy,cardw,ch, fill=PANEL, line=LINE, lw=0.75)
    text(s, cx+0.22, cy, cardw-0.4, ch, [{"ls":1.18,"runs":[R(b, 12.5, GRAY)]}], anchor=MSO_ANCHOR.MIDDLE)
box(s, ML, 5.85, CW, 0.7, fill=TEAL, line=None)
text(s, ML, 5.85, CW, 0.7, [{"align":PP_ALIGN.CENTER,"runs":[R("CowTalk은 바로 이 간극 — 알람과 행동 사이 — 를 메웁니다.", 15, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 5. 해법 (무엇을) — 4층 구조
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "해법  ·  CowTalk이란", "센서 위에 공공데이터·AI·역할별 행동을 얹는 ‘축산 운영체제’입니다.", phase="무엇을")
layers = [
 ("L4", "인텔리전스 루프", "피드백 → 정확도 추적 → 프롬프트 개선 (데이터가 쌓일수록 똑똑해진다)", TEALD),
 ("L3", "역할별 행동", "농가 · 수의사 · 방역관 · 행정관 — 같은 데이터를 각자 관점·액션플랜으로", INK2),
 ("L2", "CowTalk AI", "Claude 해석 엔진이 맥락을 읽고 역할별 ‘할 일’ 생성 (룰엔진 fallback)", TEAL),
 ("L1", "데이터 융합", "smaXtec 센서 + 이력제 · 혈통 · 등급 · DHI · KAHIS · 기상 공공데이터", INK),
]
ly=2.2; lh=0.84; lg=0.12
for i,(tag,tt,desc,col) in enumerate(layers):
    yy = ly + i*(lh+lg)
    box(s, ML, yy, CW, lh, fill=col, line=None)
    box(s, ML, yy, 1.05, lh, fill=WHITE, line=None)
    text(s, ML, yy, 1.05, lh, [{"align":PP_ALIGN.CENTER,"runs":[R(tag, 17, col, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, ML+1.3, yy, 3.2, lh, [{"runs":[R(tt, 16, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, ML+4.5, yy, CW-4.7, lh, [{"ls":1.05,"runs":[R(desc, 12.5, RGBColor(0xDD,0xE7,0xEA))]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, ML, 6.18, CW, 0.5,
     [{"runs":[R("smaXtec을 복제하지 않습니다 — ", 14, GRAY), R("그 위에 3개 레이어(공공데이터·AI·다중역할)를 더합니다.", 14, INK, True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 6. 현황 데이터 (무엇을 — 실운영 증거)
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "현황  ·  지금 작동 중", "시연이 아닙니다 — 지금 199개 농장·10,886두가 실시간으로 돌아갑니다.", phase="무엇을")
stats = [("199","개 농장","실운영"),("10,886","두","위내센서 모니터링"),("730,491","건","누적 분석 이벤트"),("99.4%","센서 가동","실시간 수집")]
sw=(CW-3*0.3)/4; sy=2.2; sh=1.5
for i,(big,unit,sub) in enumerate(stats):
    cx=ML+i*(sw+0.3)
    box(s, cx,sy,sw,sh, fill=PANEL, line=LINE, lw=0.75)
    text(s, cx, sy+0.22, sw, 0.7, [{"align":PP_ALIGN.CENTER,"runs":[R(big, 33, TEAL, True)]}])
    text(s, cx, sy+0.92, sw, 0.3, [{"align":PP_ALIGN.CENTER,"runs":[R(unit, 13, INK, True)]}])
    text(s, cx, sy+1.18, sw, 0.3, [{"align":PP_ALIGN.CENTER,"runs":[R(sub, 10.5, MUTE)]}])
# 경기 강조 + 시도 바차트
by=4.05
box(s, ML, by, CW, 2.0, fill=WHITE, line=LINE, lw=1.0)
box(s, ML, by, 0.08, 2.0, fill=TEAL, line=None)
text(s, ML+0.35, by+0.24, 5.5, 0.7, [{"ls":1.05,"runs":[R("경기 52곳 = 수도권 최대 · 전국 2위 거점", 16, INK, True)]}])
text(s, ML+0.35, by+1.0, 5.4, 0.9,
     [{"ls":1.25,"runs":[R("단일 광역으로 데이터·노하우가 가장 집약될 수 있는 곳.\n데이터가 가장 많은 곳이 ", 12, GRAY), R("표준", 12, TEAL, True), R("을 만든다.", 12, GRAY)]}])
box(s, ML+6.05, by+0.3, 0.015, 1.4, fill=LINE, line=None)
# 바차트 (충북77/경기52/충남27) — 전국 분포
bars=[("충북",77,GRAY,"77곳"),("경기",52,TEAL,"52곳 · 수도권 최대"),("충남",27,GRAY,"27곳")]
name_x=ML+6.3; bar_x=ML+7.05; bw_max=3.0; bh=0.3; bgap=0.24; bsy=by+0.55
text(s, name_x, by+0.22, 4.5, 0.3, [{"runs":[R("시도별 농장 수 (전국 199)", 11, MUTE, True)]}])
for i,(nm,val,col,lab) in enumerate(bars):
    yy=bsy+i*(bh+bgap)
    text(s, name_x, yy, 0.72, bh, [{"runs":[R(nm, 12, INK, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    box(s, bar_x, yy, bw_max*val/77.0, bh, fill=col, line=None)
    text(s, bar_x+bw_max*val/77.0+0.1, yy, 2.4, bh, [{"runs":[R(lab, 11, (TEALD if col==TEAL else MUTE), col==TEAL)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 7. 어떻게 — 알람→행동 완결
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "작동 방식  ·  알람을 행동으로", "알람 → 판단 → 추천 → 행동 → 기록을 한 화면에서 끝냅니다.", phase="어떻게")
steps=[("알람","위내센서 발열·발정 감지"),("판단","AI가 맥락 해석"),("추천","수정적기·정액·치료 제시"),("행동","담당자에 즉시 배정"),("기록","결과 피드백 → 학습")]
sw=(CW-4*0.28)/5; sy=2.4; sh=1.55
for i,(tt,desc) in enumerate(steps):
    cx=ML+i*(sw+0.28)
    col = TEAL if i==0 else INK if i==4 else INK2
    box(s, cx,sy,sw,sh, fill=(PANEL2 if i in(0,4) else PANEL), line=LINE, lw=0.75)
    text(s, cx, sy+0.2, sw, 0.4, [{"align":PP_ALIGN.CENTER,"runs":[R(f"{i+1}", 15, WHITE, True)]}])
    box(s, cx+sw/2-0.2, sy+0.22, 0.4, 0.4, fill=col, line=None, rounded=True, radius=0.5)
    text(s, cx, sy+0.22, sw, 0.4, [{"align":PP_ALIGN.CENTER,"runs":[R(f"{i+1}", 14, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx, sy+0.78, sw, 0.35, [{"align":PP_ALIGN.CENTER,"runs":[R(tt, 15, INK, True)]}])
    text(s, cx+0.12, sy+1.12, sw-0.24, 0.4, [{"align":PP_ALIGN.CENTER,"ls":1.05,"runs":[R(desc, 10.5, GRAY)]}])
    if i<4:
        text(s, cx+sw-0.02, sy, 0.3, sh, [{"align":PP_ALIGN.CENTER,"runs":[R("›", 22, MUTE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
box(s, ML, 4.5, CW, 1.55, fill=INK, line=None)
text(s, ML+0.4, 4.7, CW-0.8, 0.4, [{"runs":[R("팅커벨 AI — 말 한마디로 끝낸다", 15, MINT, True)]}])
text(s, ML+0.4, 5.18, CW-0.8, 0.8,
     [{"ls":1.25,"runs":[R("“경기도 발열 현황 보여줘”  →  ", 14, WHITE),
       R("3초", 16, MINT, True),
       R("에 경기 농장 발열 현황이 한 화면에.  엑셀·전화가 대화 한 줄로.", 14, WHITE)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 8. 그래서 결과 ① 방역
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "그래서 결과는  ·  방역 주권", "역학조사가 며칠에서 30초로 — 1회 차단이 시스템값을 회수합니다.", phase="그래서 결과는")
# before/after
half=(CW-0.4)/2
box(s, ML, 2.25, half, 2.3, fill=PANEL, line=LINE, lw=1.0)
text(s, ML+0.3, 2.45, half-0.6, 0.4, [{"runs":[R("AS-IS  기존 방역", 13, GRAY, True)]}])
text(s, ML+0.3, 2.95, half-0.6, 1.5,
     [{"ls":1.35,"runs":[R("· 전화·수기 엑셀 역학조사\n· 접촉 농장 파악에 며칠\n· 확산은 늘 사후 대응", 14, INK)]}])
box(s, ML+half+0.4, 2.25, half, 2.3, fill=PANEL2, line=TEAL, lw=1.2)
text(s, ML+half+0.7, 2.45, half-0.6, 0.4, [{"runs":[R("TO-BE  CowTalk 방역", 13, TEALD, True)]}])
text(s, ML+half+0.7, 2.95, half-0.6, 1.5,
     [{"ls":1.35,"runs":[R("· 접촉망 자동 추적 ", 14, INK), R("30초", 14, TEALD, True),
       R("\n· 증상 전 위내센서 선제 감지\n· 199농장 전국 통합 실시간 감시", 14, INK)]}])
# 결과 배너
box(s, ML, 4.85, CW, 1.25, fill=INK, line=None)
text(s, ML+0.45, 4.85, 5.9, 1.25,
     [{"ls":1.15,"runs":[R("전염병 1회 차단 = ", 15, WHITE), R("수십~수백억", 24, MINT, True),
       R("\n보상·살처분 비용 회피", 13, RGBColor(0xC7,0xD6,0xDE))]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, EW-MR-4.6, 4.85, 4.2, 1.25,
     [{"align":PP_ALIGN.RIGHT,"ls":1.15,"runs":[R("“한 번만 막아도\n시스템값을 회수한다.”", 15, WHITE, True, it=True)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 9. 그래서 결과 ② 농가 경제
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "그래서 결과는  ·  농가 소득", "번식·질병을 데이터로 잡아 젖소 두당 연 24만원의 생산성을 만듭니다.", phase="그래서 결과는")
comps=[("번식효율 향상",10,TEAL),("질병 조기대응",6,INK2),("산유량 향상",5,TEALD),("노동 절감",3,GRAY)]
total=24.0; barx=ML; bary=2.45; barW=CW; barH=0.95
acc=0
for nm,val,col in comps:
    w=barW*val/total
    box(s, barx+acc, bary, w, barH, fill=col, line=WHITE, lw=1.5)
    text(s, barx+acc, bary, w, barH, [{"align":PP_ALIGN.CENTER,"ls":1.05,"runs":[R(f"{val}만\n", 17, WHITE, True), R(nm, 10.5, WHITE)]}], anchor=MSO_ANCHOR.MIDDLE)
    acc+=w
text(s, ML, bary+barH+0.2, CW, 0.4, [{"runs":[R("젖소 ", 14, GRAY), R("두당 24만원 / 년", 16, INK, True), R("  생산성 편익 (4개 항목 합산)", 13, GRAY)]}])
# 합계 배너
box(s, ML, 4.4, CW, 1.6, fill=PANEL2, line=None)
text(s, ML+0.45, 4.4, 8.5, 1.6,
     [{"runs":[R("경기 젖소 ", 16, INK), R("100,000두", 16, TEALD, True), R(" 적용 시", 16, INK)]},
      {"sb":6,"runs":[R("농가 직접 편익  ", 14, GRAY), R("240억 / 년", 30, TEALD, True)]}], anchor=MSO_ANCHOR.MIDDLE)
text(s, EW-MR-3.6, 4.4, 3.2, 1.6,
     [{"align":PP_ALIGN.RIGHT,"ls":1.2,"runs":[R("알람이 행동으로\n완결될 때\n생기는 돈", 13, GRAY)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 10. 그래서 결과 ③ 사회·환경
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "그래서 결과는  ·  사회·환경", "적게 키우고 더 잘 짠다 — 육성우 감축으로 사료·분뇨·탄소를 동시에 줄입니다.", phase="그래서 결과는")
box(s, ML, 2.2, CW, 0.78, fill=INK, line=None)
text(s, ML+0.35, 2.2, CW-0.7, 0.78,
     [{"runs":[R("핵심: 사육두수 ↓ 이지만 우유 총생산량은 유지.   ", 14, WHITE, True),
       R("경기 육성우 ~80,000두 중 36,000두(45%) 감축 가능.", 13, RGBColor(0xC7,0xD6,0xDE))]}], anchor=MSO_ANCHOR.MIDDLE)
cards=[("사료비 절감","432억","/ 년","농가 직접 비용↓",TEAL),
       ("분뇨 배출 감소","28.8만톤","처리비 86억","악취·민원 완화",INK2),
       ("탄소 저감","10.8만톤","CO₂eq · 108억","탄소중립 기여",TEALD)]
cw=(CW-2*0.4)/3; cy=3.25; ch=2.0
for i,(tt,big,sub,note,col) in enumerate(cards):
    cx=ML+i*(cw+0.4)
    box(s, cx,cy,cw,ch, fill=WHITE, line=LINE, lw=1.0)
    box(s, cx,cy,cw,0.08, fill=col, line=None)
    text(s, cx+0.28, cy+0.28, cw-0.5, 0.35, [{"runs":[R(tt, 14, INK, True)]}])
    text(s, cx+0.28, cy+0.74, cw-0.5, 0.6, [{"runs":[R(big, 28, col, True)]}])
    text(s, cx+0.28, cy+1.36, cw-0.5, 0.3, [{"runs":[R(sub, 12, GRAY)]}])
    text(s, cx+0.28, cy+1.64, cw-0.5, 0.3, [{"runs":[R(note, 11, MUTE)]}])
box(s, ML, 5.5, CW, 0.62, fill=AMBER, line=None)
text(s, ML, 5.5, CW, 0.62, [{"align":PP_ALIGN.CENTER,"runs":[R("사회·환경 편익 소계  ≈  626억 / 년", 16, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
footer(s)

# ══════════════════════════════════════════════════════════════════
# 11. 예산·ROI
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "그래서 결과는  ·  예산·ROI", "연 180억은 비용이 아니라 보험료 — 통합 편익 896억, BCR ≈ 5.0.", phase="그래서 결과는")
text(s, ML, 1.98, CW, 0.32, [{"runs":[R("전제  ", 11.5, MUTE, True), R("경기 젖소 100,000두 · 1,500농장 · 중계기 500만/농가 · 두당 월 15,000원(센서+CowTalk)", 11.5, GRAY)]}])
# 좌: 비용/편익 막대
chx=ML; chy=2.5; chMaxH=2.55; baseY=chy+chMaxH
def vbar(cx, val, maxv, segs, label, total_lab):
    h=chMaxH*val/maxv
    yy=baseY-h
    acc=0
    for sv,scol,sl in segs:
        sh=chMaxH*sv/maxv
        box(s, cx, baseY-acc-sh, 1.7, sh, fill=scol, line=WHITE, lw=1.2)
        if sh>0.4:
            text(s, cx, baseY-acc-sh, 1.7, sh, [{"align":PP_ALIGN.CENTER,"runs":[R(sl, 11, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
        acc+=sh
    text(s, cx-0.3, baseY+0.08, 2.3, 0.3, [{"align":PP_ALIGN.CENTER,"runs":[R(label, 12, INK, True)]}])
    text(s, cx-0.3, baseY+0.36, 2.3, 0.4, [{"align":PP_ALIGN.CENTER,"runs":[R(total_lab, 18, INK, True)]}])
maxv=900
vbar(ML+0.5, 180, maxv, [(180, GRAY, "비용 180")], "연 비용", "180억")
vbar(ML+3.2, 896, maxv, [(240, TEAL,"농가 240"),(30, INK2,"방역 30"),(626, TEALD,"사회·환경 626")], "연 편익", "896억")
# 화살표/배수
text(s, ML+5.2, chy+0.8, 1.4, 1.0, [{"align":PP_ALIGN.CENTER,"runs":[R("× 5.0", 28, AMBER, True)]}], anchor=MSO_ANCHOR.MIDDLE)
# 우: 비용 추이 + ROI 배지
rx=ML+6.7
box(s, rx, chy, CW-(rx-ML), 1.55, fill=PANEL, line=LINE, lw=0.75)
text(s, rx+0.3, chy+0.18, 4.5, 0.3, [{"runs":[R("도입 비용 (행정 예산 관점)", 12.5, INK, True)]}])
rows=[("1년차 (구축 포함)","255억"),("2년차~ 매년","180억"),("5년 누적","975억")]
for i,(a,b) in enumerate(rows):
    yy=chy+0.6+i*0.3
    text(s, rx+0.3, yy, 3.2, 0.3, [{"runs":[R(a, 12, GRAY)]}])
    text(s, rx+3.6, yy, 1.6, 0.3, [{"align":PP_ALIGN.RIGHT,"runs":[R(b, 13, INK, True)]}])
badges=[("통합 BCR","≈ 5.0"),("회수기간","~2.8년"),("5년 순편익(농가)","+195억")]
bw=(CW-(rx-ML)-2*0.2)/3
for i,(t,v) in enumerate(badges):
    cx=rx+i*(bw+0.2)
    box(s, cx, chy+1.75, bw, 1.15, fill=INK, line=None)
    text(s, cx, chy+1.95, bw, 0.5, [{"align":PP_ALIGN.CENTER,"runs":[R(v, 20, MINT, True)]}])
    text(s, cx, chy+2.5, bw, 0.3, [{"align":PP_ALIGN.CENTER,"runs":[R(t, 10.5, RGBColor(0xC7,0xD6,0xDE))]}])
text(s, ML, 6.15, CW, 0.5, [{"runs":[R("“농가 생산성으로 3년 안에 회수, 사료·분뇨·탄소 동시 절감, 전염병 한 번만 막으면 그 자체로 끝.”", 13, INK, it=True)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 12. 실행 로드맵
# ══════════════════════════════════════════════════════════════════
s = slide()
header(s, "실행  ·  어떻게 확산하나", "1개 농장에서 전국·수출까지 — 검증된 구조를 단계로 확장합니다.", phase="실행")
phases=[("지금","199 농장 · 경기 52","실운영 검증 완료",TEAL),
        ("1단계","경기 시범 (시군+공수의사)","실증·효과 측정",INK2),
        ("2단계","경기 전역 1,500 농장","도 단위 표준화",INK2),
        ("3단계","전국 표준","국가 방역 고도화",INK2),
        ("4단계","해외 수출","i18n·국가별 어댑터",TEALD)]
pw=(CW-4*0.25)/5; py=2.35; ph=1.7
hline(s, ML+pw/2, py+0.35, CW-pw, color=LINE, weight=1.5)
for i,(st,tt,desc,col) in enumerate(phases):
    cx=ML+i*(pw+0.25)
    box(s, cx+pw/2-0.11, py+0.24, 0.22, 0.22, fill=col, line=WHITE, lw=1.5, rounded=True, radius=0.5)
    box(s, cx, py+0.62, pw, ph-0.62, fill=(PANEL2 if i==0 else PANEL), line=LINE, lw=0.75)
    text(s, cx, py+0.78, pw, 0.3, [{"align":PP_ALIGN.CENTER,"runs":[R(st, 12, col, True)]}])
    text(s, cx+0.1, py+1.12, pw-0.2, 0.5, [{"align":PP_ALIGN.CENTER,"ls":1.05,"runs":[R(tt, 11.5, INK, True)]}])
    text(s, cx+0.1, py+1.62, pw-0.2, 0.3, [{"align":PP_ALIGN.CENTER,"runs":[R(desc, 10, MUTE)]}])
# 경기도 요청
text(s, ML, 4.55, CW, 0.35, [{"runs":[R("경기도에 드리는 요청", 14, INK, True)]}])
asks=[("시범사업","경기 일부 시군 + 공수의사 대상 실증"),
      ("데이터 연계","KAHIS · 이력제 · 도 방역 데이터 연동 협조"),
      ("확산 로드맵","경기 전역 → 전국 표준화 공동 수립")]
aw=(CW-2*0.4)/3; ay=4.95; ah=1.05
for i,(t,d) in enumerate(asks):
    cx=ML+i*(aw+0.4)
    box(s, cx,ay,aw,ah, fill=WHITE, line=TEAL, lw=1.0)
    text(s, cx+0.25, ay+0.18, aw-0.5, 0.35, [{"runs":[R(f"{i+1}. {t}", 13.5, TEALD, True)]}])
    text(s, cx+0.25, ay+0.56, aw-0.5, 0.4, [{"ls":1.1,"runs":[R(d, 11.5, GRAY)]}])
footer(s)

# ══════════════════════════════════════════════════════════════════
# 13. 제언 closing
# ══════════════════════════════════════════════════════════════════
s = slide()
box(s, 0,0, EW, EH, fill=INK, line=None)
box(s, 0,0, 0.22, EH, fill=TEAL, line=None)
text(s, 1.0, 1.5, 11.4, 0.4, [{"runs":[R("제언", 13, MINT, True)]}])
text(s, 1.0, 2.0, 11.4, 1.8,
     [{"ls":1.08,"runs":[R("경기도가 대한민국 축산\n디지털 행정의 첫 모델이 됩니다.", 36, WHITE, True)]}])
hline(s, 1.0, 4.2, 4.2, color=TEAL, weight=2.2)
text(s, 1.0, 4.45, 11.0, 0.8,
     [{"ls":1.3,"runs":[R("데이터가 가장 많은 곳이 표준을 만든다.  경기 52곳, 수도권 최대 거점에서\n검증된 구조를 전국·세계로 확장합니다.", 15, RGBColor(0xC7,0xD6,0xDE))]}])
mini=[("방역 주권","전염병 1회 = 수백억 회피"),("농가·사회편익","896억 / 년 (BCR ≈ 5.0)"),("확장","경기 → 전국 → 수출")]
mw=(11.4-2*0.4)/3
for i,(t,d) in enumerate(mini):
    cx=1.0+i*(mw+0.4)
    box(s, cx, 5.6, mw, 0.95, fill=RGBColor(0x12,0x2E,0x47), line=RGBColor(0x24,0x44,0x5E), lw=0.75)
    text(s, cx+0.25, 5.75, mw-0.5, 0.3, [{"runs":[R(t, 13, MINT, True)]}])
    text(s, cx+0.25, 6.1, mw-0.5, 0.35, [{"runs":[R(d, 12, WHITE)]}])
text(s, 1.0, 6.9, 11.4, 0.3, [{"runs":[R("CowTalk v5.0  ·  D2O Corp  ·  함께 만드는 축산 디지털 주권", 11, MUTE)]}])

# ── 발표자 노트 (15분 운영, 무엇을→어떻게→그래서 결과는) ──
NOTES = [
 "[0:30] 인사. ‘오늘 15분 안에 — 무엇을 만들었고, 어떻게 작동하며, 그래서 경기도에 어떤 결과가 오는지 — 숫자로 말씀드리겠습니다.’",
 "[1:30] 답부터. ‘결론은 하나입니다 — 경기도가 선도해야 합니다.’ 방역·농가·사회환경 3축. 연 180억 투자 → 편익 896억(BCR≈5.0). 이 한 장이 전체 요약입니다.",
 "[1:00] 왜 경기인가. 대한민국 낙농의 중심은 경기다 — 젖소 사육 전국 약 41%, 수도권 인구 50%+ 로 최대 소비지도 경기. 생산·소비를 동시에 가진 유일 광역. 그래서 축우(젖소·한우)부터 디지털 대전환을 시작해 양돈·가금 등 타 축종으로 확산. 경기가 표준화하면 그것이 곧 대한민국 표준.",
 "[1:00] 현황. 세계는 축산을 데이터로 운영하는데 우리 방역·행정은 아직 전화와 엑셀. 데이터는 쌓이는데 판단·행동으로 안 이어진다.",
 "[1:00] 핵심 공백. smaXtec 센서·알람은 세계 최고지만 ‘알람→행동’의 다리가 없다. 농가·방역관이 ‘그래서 뭘?’을 혼자 판단. CowTalk이 이 간극을 메운다.",
 "[1:30] 무엇을. CowTalk=축산 운영체제. smaXtec 복제가 아니라 그 위에 공공데이터·AI·역할별 행동 3개 층을 얹는다. 4층 구조를 짧게 짚기.",
 "[1:30] 시연이 아니다. 지금 199농장·10,886두가 실시간으로 돈다. 경기 52곳=수도권 최대·전국 2위(전국 1위는 충북 77 — 사실대로). 데이터 가장 많은 곳이 표준을 만든다.",
 "[1:30] 어떻게. 알람→판단→추천→행동→기록을 한 화면에서 완결. 팅커벨에 ‘경기도 발열 현황 보여줘’ → 3초. (가능하면 실제 시연 1회)",
 "[1:30] 그래서 결과①. 역학조사 며칠→30초, 증상 전 선제 감지. 전염병 1회 차단=수십~수백억 회피. ‘한 번만 막아도 시스템값 회수.’",
 "[1:00] 그래서 결과②. 번식·질병을 데이터로 잡아 두당 24만원/년. 경기 10만두면 240억/년 농가 직접 편익.",
 "[1:30] 그래서 결과③. 핵심 메시지 ‘적게 키우고 더 잘 짠다.’ 육성우 36,000두 감축 → 사료 432억+분뇨 86억+탄소 108억 = 626억/년. 탄소중립·민원까지.",
 "[2:00] 핵심 슬라이드. 천천히. 연 180억은 비용이 아니라 보험료. 편익 896억(농가240+방역30+사회환경626)=BCR≈5.0, 회수 ~2.8년. 한 줄로 마무리.",
 "[1:00] 어떻게 확산. 199농장(경기52)에서 검증 → 경기 시범 → 경기 전역 1,500 → 전국 → 수출. 경기도 요청 3가지(시범사업·데이터 연계·로드맵 공동수립).",
 "[0:30] 클로징. ‘경기도가 대한민국 축산 디지털 행정의 첫 모델이 됩니다.’ 데이터가 가장 많은 곳이 표준을 만든다. 감사합니다.",
]
for sl, nt in zip(prs.slides, NOTES):
    sl.notes_slide.notes_text_frame.text = nt

out="/tmp/claude-0/-home-user-cowtalk/3e10b74e-cb9f-572d-bcc2-894bf6a7f2ee/scratchpad/CowTalk_경기인수위_브리핑.pptx"
prs.save(out)
print("SAVED", out, "| slides:", len(prs.slides._sldIdLst))
