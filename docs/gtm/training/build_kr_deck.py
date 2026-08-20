# -*- coding: utf-8 -*-
"""국문판 빌더 — 원본 35장을 번역하고, CowTalk AI 28장을 국문으로 이어 붙인다.
한글은 eastAsia 폰트를 '맑은 고딕'으로 강제한다 (원본 테마는 Calibri라 한글 글리프가 없다)."""
import os, sys, zipfile, html, shutil
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from tr import TR, FOOTER_EN, FOOTER_KR

UN, EMU = 'un', 914400
KRFONT = '맑은 고딕'
def E(v): return str(int(round(v * EMU)))
def esc(s): return html.escape(s, quote=False)

BG_LIGHT, BG_DARK = 'F5F3EA', '0A1B17'
INK, MUTED = '10251F', '66746F'
TEAL, GREEN = '15A98C', '65D44A'
CORAL, BLUE, AMBER = 'E96155', '4C8DD9', 'F6B846'
WHITE, PANEL = 'FFFFFF', '1BAF91'
DK_SUB, DK_FOOT = 'C9D8D2', 'BFD1CA'
CALLOUT_BG = 'DFF4E8'
ML, CW = 0.75, 11.83

_uid = [2000]
def nid():
    _uid[0] += 1; return _uid[0]

def _font():
    return f'<a:latin typeface="Calibri"/><a:ea typeface="{KRFONT}"/><a:cs typeface="{KRFONT}"/>'

def sp(x, y, w, h, geom='rect', fill=None, runs=None, align='l'):
    i = nid()
    f = f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>' if fill else '<a:noFill/>'
    body = '<a:p><a:endParaRPr lang="ko-KR" altLang="en-US"/></a:p>'
    if runs:
        paras = []
        for para in runs:
            rs = ''
            for r in para:
                sz = int(r.get('sz', 16.5) * 100); col = r.get('c', INK)
                b = ' b="1"' if r.get('b') else ''
                rs += (f'<a:r><a:rPr lang="ko-KR" altLang="en-US" sz="{sz}"{b}>'
                       f'<a:solidFill><a:srgbClr val="{col}"/></a:solidFill>{_font()}</a:rPr>'
                       f'<a:t>{esc(r["t"])}</a:t></a:r>')
            sz = int(para[0].get('sz', 16.5) * 100)
            paras.append(f'<a:p><a:pPr algn="{align}"><a:buNone/>'
                         f'<a:defRPr sz="{sz}">{_font()}</a:defRPr></a:pPr>{rs}</a:p>')
        body = ''.join(paras)
    return (f'<p:sp><p:nvSpPr><p:cNvPr id="{i}" name="k{i}"/>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>'
            f'<p:spPr><a:xfrm><a:off x="{E(x)}" y="{E(y)}"/><a:ext cx="{E(w)}" cy="{E(h)}"/></a:xfrm>'
            f'<a:prstGeom prst="{geom}"><a:avLst/></a:prstGeom>{f}'
            f'<a:ln w="0"><a:noFill/><a:prstDash val="solid"/></a:ln></p:spPr>'
            f'<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0"/>'
            f'<a:lstStyle/>{body}</p:txBody></p:sp>')

def text(x, y, w, h, t, sz=16.5, c=INK, b=False, align='l'):
    return sp(x, y, w, h, runs=[[{'t': ln, 'sz': sz, 'c': c, 'b': b}] for ln in t.split('\n')], align=align)

def chrome(page, dark=False):
    return [text(ML, 7.09, 7.29, 0.19, FOOTER_KR, 8.25, DK_FOOT if dark else MUTED),
            text(12.29, 7.09, 0.29, 0.19, f'{page:02d}', 8.25, WHITE if dark else INK, b=True, align='r')]

def eyebrow_title(e, t):
    return [text(ML, 0.44, 5.42, 0.25, e, 10.5, TEAL, b=True),
            text(ML, 0.81, CW, 1.00, t, 30, INK, b=True)]

def callout(t):
    return [sp(ML, 5.31, CW, 0.77, geom='roundRect', fill=CALLOUT_BG),
            text(1.06, 5.54, CW - 0.62, 0.33, t, 16.5, INK, b=True)]

def bullets(items, x=0.79, tx=1.08, w=11.50, y0=2.06, pitch=0.73, sz=16.5, h=0.69):
    o = []
    for k, it in enumerate(items):
        yy = y0 + k * pitch
        o += [sp(x, yy + 0.08, 0.125, 0.125, geom='ellipse', fill=GREEN),
              text(tx, yy, w, h, it, sz, INK)]
    return o

def card(x, y, w, h, accent, label, title, body, body_sz=12.75):
    o = [sp(x, y, w, h, geom='roundRect', fill=WHITE), sp(x, y, w, 0.12, fill=accent),
         text(x + 0.29, y + 0.38, w - 0.58, 0.30, label, 12.75, accent, b=True),
         text(x + 0.29, y + 0.84, w - 0.58, 0.70, title, 16.5, INK, b=True)]
    if body: o.append(text(x + 0.29, y + 1.68, w - 0.58, h - 1.95, body, body_sz, MUTED))
    return o

def row_x(n, gap=0.34):
    w = (CW - gap * (n - 1)) / n
    return w, [ML + k * (w + gap) for k in range(n)]

def stat(x, y, w, value, label, color=CORAL):
    return [text(x, y, w, 0.67, value, 36, color, b=True),
            text(x, y + 0.65, w, 0.75, label, 12.75, MUTED)]

def steps(items, y=2.20, h=2.60, gap=0.22, sz=13.5, colors=None):
    n = len(items); w = (CW - gap * (n - 1)) / n; o = []
    for k, (num, head, desc) in enumerate(items):
        x = ML + k * (w + gap); cc = (colors[k] if colors else None) or TEAL
        o += [sp(x, y, w, h, geom='roundRect', fill=WHITE),
              sp(x + 0.26, y + 0.28, 0.42, 0.42, geom='ellipse', fill=cc),
              text(x + 0.26, y + 0.34, 0.42, 0.30, num, 13.5, WHITE, b=True, align='ctr'),
              text(x + 0.26, y + 0.92, w - 0.52, 0.62, head, 15, INK, b=True),
              text(x + 0.26, y + 1.66, w - 0.52, h - 1.92, desc, sz, MUTED)]
    return o

def layer_rows(rows, y0=1.96, h=0.70, gap=0.09):
    o = []
    for k, (tag, name, desc, col) in enumerate(rows):
        yy = y0 + k * (h + gap)
        o += [sp(ML, yy, CW, h, geom='roundRect', fill=WHITE), sp(ML, yy, 0.12, h, fill=col),
              text(ML + 0.34, yy + 0.12, 1.55, 0.26, tag, 11.25, col, b=True),
              text(ML + 2.00, yy + 0.09, 3.55, 0.32, name, 15.75, INK, b=True),
              text(ML + 5.75, yy + 0.11, CW - 6.05, 0.50, desc, 13.5, MUTED)]
    return o

def two_col(left, right, y=2.06, h=3.05):
    w, xs = row_x(2, gap=0.43)
    n = max(len(left[3]), len(right[3]))
    pitch, ih = (0.58, 0.54) if n <= 3 else (0.44, 0.40)
    o = []
    for (x, (accent, label, title, items)) in zip(xs, (left, right)):
        o += [sp(x, y, w, h, geom='roundRect', fill=WHITE), sp(x, y, w, 0.12, fill=accent),
              text(x + 0.30, y + 0.36, w - 0.60, 0.28, label, 12.75, accent, b=True),
              text(x + 0.30, y + 0.76, w - 0.60, 0.44, title, 17.25, INK, b=True)]
        for k, it in enumerate(items):
            yy = y + 1.30 + k * pitch
            o += [sp(x + 0.32, yy + 0.06, 0.11, 0.11, geom='ellipse', fill=accent),
                  text(x + 0.58, yy, w - 0.92, ih, it, 13.5, MUTED)]
    return o

def divider(num, title, sub, page):
    o = [text(ML, 0.77, 2.29, 1.35, f'{num:02d}', 69, GREEN, b=True),
         text(ML, 2.66, 9.27, 1.04, title, 40.5, WHITE, b=True),
         text(ML, 3.96, 8.12, 1.10, sub, 17.25, DK_SUB),
         sp(10.83, 0.0, 2.50, 7.50, fill=PANEL)]
    for (x, y, d, c) in [(11.35,5.89,0.19,WHITE),(11.62,5.57,0.27,GREEN),(11.90,5.26,0.35,WHITE),
                         (12.17,4.95,0.44,GREEN),(12.44,4.64,0.52,WHITE)]:
        o.append(sp(x, y, d, d, geom='ellipse', fill=c))
    return o + chrome(page, dark=True), True

def content(e, t, blocks, page):
    return eyebrow_title(e, t) + blocks + chrome(page), False

SLIDES = []
def add(b, n): SLIDES.append((b, n))
P = [36]
def nxt():
    p = P[0]; P[0] += 1; return p

# ── 36 연결
add(content('센서에서 시스템으로', '알람이 끝나는 곳에서 CowTalk이 시작합니다', two_col(
    (TEAL, '센서가 주는 것', '신뢰할 수 있는 사건',
     ['소의 몸속에서 잡아낸 생리적 이탈', '시각이 찍히고 제조사가 검증한 신호', '어느 목장에서나 같은 품질로']),
    (CORAL, '그다음에 남는 것', '사람이 내려야 할 결정',
     ['알람 열한 개 중 어느 소부터 봐야 하는가', '그 소의 이력은 — 치료, 산차, 착유일수는', '누가 무엇으로 조치하고, 실제로 했는가'])) +
    callout('알람과 행동 사이의 이 간극에서 가치의 대부분이 사라집니다.'), nxt()),
    '이 장이 교육 전체의 경첩이다. 앞의 35장은 센서를 다뤘고, 여기서부터는 센서 이벤트를 결정·행동·기록으로 바꾸는 층을 다룬다. '
    'CowTalk을 "더 좋은 센서"로 소개하지 말 것 — 층이 다르다.')

# ── 37 §07
add(divider(7, '센서 위의 레이어',
    'CowTalk은 센서 이벤트 위에 맥락과 해석, 역할별 행동을 더합니다. 이벤트 자체를 다시 판단하지는 않습니다.', nxt()),
    '경계를 소리 내어 말한다. 제조사가 검증한 이벤트는 신뢰한다. 그걸 다시 채점하면 정확도만 떨어진다. 우리 일은 이벤트 이후부터다.')

# ── 38 4층
add(content('구조', '네 개의 층이 신호를 결정으로 바꿉니다', layer_rows([
    ('LAYER 4', '학습 루프', '전문가 레이블·알람 통계·치료 결과가 다음 판단을 개선한다', AMBER),
    ('LAYER 3', '역할별 서빙', '같은 데이터가 농가·수의사·방역관·행정에 각기 다르게 도달한다', BLUE),
    ('LAYER 2', '해석', '이벤트에 맥락을 더한다. 모든 판단은 근거와 출처를 함께 남긴다', TEAL),
    ('LAYER 1', '데이터 통합', '센서 이벤트·국가 기록·농장 이력·기상을 한 마리에 모은다', GREEN),
]) + callout('센서 이벤트는 신뢰하고 재판단하지 않습니다. 센서가 알 수 없는 것을 더할 뿐입니다.'), nxt()),
    '아래에서 위로 읽는다. 데이터가 들어오는 곳은 1층뿐이고, 2층은 그것을 우회하지 못한다. '
    '어시스턴트가 기억만으로 답할 수 없고 반드시 데이터를 조회해야 한다는 제약이 감사 가능성을 만든다.')

# ── 39 네 개의 번호
w4, x4 = row_x(4)
add(content('맥락', '소 한 마리에 네 개의 번호가 따로 살고 있었습니다',
    card(x4[0], 2.08, w4, 3.05, CORAL, '이력제 번호', '국가 등록 번호', '출생 2주 안에 국가가 부여한다. 공공 기록을 여는 열쇠다.') +
    card(x4[1], 2.08, w4, 3.05, TEAL, '목장 관리번호', '현장에서 부르는 이름', '실제로 목장에서 부르는 번호. 모든 대화가 이 번호로 오간다.') +
    card(x4[2], 2.08, w4, 3.05, BLUE, '혈통등록번호', '혈통과 유전능력', '근교계수 확인과 정액 선택에 필요하다.') +
    card(x4[3], 2.08, w4, 3.05, AMBER, '센서 시리얼', '장치 식별자', '몸속의 볼루스. 생리 신호를 위의 모든 것과 잇는다.') +
    callout('이 넷을 잇는 순간 "열이 난다"가 우선순위와 손실 금액으로 바뀝니다.'), nxt()),
    '대부분의 나라에서 이 네 번호는 서로 만나지 않는 네 시스템에 흩어져 있다. '
    '"지금 개체 한 마리의 치료 이력을 확인하는 데 얼마나 걸립니까"를 청중에게 물어보라. 그 대답이 이 장의 가치다.')

# ── 40 경계
add(content('경계', '어시스턴트가 하는 일과 하지 않는 일', two_col(
    (TEAL, '하는 일', '정리하고 근거를 붙인다',
     ['그 개체의 맥락을 한자리에 모은다', '가능한 원인을 근거와 함께 순위로 제시한다',
      '다음에 무엇을 확인할지와 시점을 제안한다', '사람이 조치한 뒤 기록을 대신 쓴다']),
    (CORAL, '하지 않는 일', '진단하거나 처방하지 않는다',
     ['진단하지 않는다 — 후보를 제시할 뿐이다', '처방하지 않는다 — 수의사의 권한이다',
      '승인 없이 실행하지 않는다', '데이터 없이 답하지 않는다 — “자료 없음”도 답이다'])) +
    callout('모든 출력은 출처를 남깁니다. 근거 없는 답변은 결과가 아니라 실패로 봅니다.'), nxt()),
    '이 장에 시간을 준다. 어느 자리에나 AI를 위험 요소로 미리 규정한 사람이 있다. 오른쪽 칸이 그 사람을 위한 것이다. '
    '최종 판단은 사람의 몫이며, 구조적으로 우회할 수 없게 만들었다고 분명히 말한다.')

# ── 41 §08
add(divider(8, '어시스턴트',
    '평소 쓰는 말로 물어보면 됩니다. 답은 그 목장의 데이터에서 나오고, 숫자마다 어디서 왔는지를 함께 보여줍니다.', nxt()),
    '통신이 되면 실연을 권한다. 안 되면 부록의 화면 캡처를 쓴다. 남겨야 할 인상은 "일반 지식이 아니라 이 목장의 데이터로 답한다"는 것이다.')

# ── 42 도구
w5, x5 = row_x(5, gap=0.22)
add(content('도구', '어시스턴트는 정해진 도구를 통해서만 움직입니다',
    [c for (x, acc, lab, ttl, bod) in [
        (x5[0], TEAL, '센서', '개체와 신호', '개체 프로필, 이벤트 이력, 체온·활동량 추이, 기상과 더위지수'),
        (x5[1], GREEN, '농장', '우군과 기록', '농장 요약, 핵심 지표, 치료 기록, 유량 입력'),
        (x5[2], BLUE, '번식', '번식 주기', '번식 통계, 수정 적기, 정액 추천, 발정동기화 처방'),
        (x5[3], AMBER, '공공데이터', '국가 기록', '이력제, 등급판정, 경락가격, 지역 방역 현황'),
        (x5[4], CORAL, '수의', '임상 지원', '감별 후보, 치료 결과 확정, 전문가 레이블 기록'),
    ] for c in card(x, 2.08, w5, 3.05, acc, lab, ttl, bod, body_sz=11.5)] +
    callout('모든 도구 호출은 감사 로그에 기록됩니다 — 누가, 언제, 무엇을 받았는지.'), nxt()),
    '어시스턴트는 없는 기능을 지어내지 못한다. 필요한 도구가 없으면 없다고 답한다. '
    '접근 권한은 역할마다 다르다 — 농가와 방역관의 도구 목록은 같지 않다. '
    '로그가 세션 단위가 아니라 호출 단위라는 점을 언급한다. 감사에서 묻는 수준이 그 정도다.')

# ── 43 근거
add(content('근거', '모든 답변은 그 숫자가 어디서 왔는지를 함께 보여줍니다', steps([
    ('1', '평소 말로 묻는다', '"423번 어때?" — 타이핑도 되고 음성도 된다'),
    ('2', '허용된 도구만 호출', '그 역할이 쓸 수 있는 도구로, 권한 안의 개체만'),
    ('3', '기록에서 값을 읽는다', 'DB의 실제 값 — 기억이나 추정이 아니다'),
    ('4', '출처와 함께 답한다', '각 수치에 출처·시각·신뢰도가 붙는다'),
], y=2.14, h=2.72) + callout('데이터가 없으면 “자료 없음”이 정답입니다. 그럴듯한 추측이 아니라.'), nxt()),
    '4단계가 전문가 사용자에게 결정적이다. 수의사는 출처가 불분명한 숫자로 움직이지 않는다. '
    '측정값이 아니라 추정값일 때 답변이 그렇다고 밝힌다는 점을 짚는다.')

# ── 44 역할
add(content('역할', '하나의 데이터, 네 개의 질문',
    card(x4[0], 2.08, w4, 3.05, GREEN, '농장주', '오늘 뭘 하지?', '오늘 할 일, 발정·분만 알림, 소 한 마리가 남기는 마진.') +
    card(x4[1], 2.08, w4, 3.05, CORAL, '수의사', '어디부터 가지?', '담당 농장 전체의 케이스 큐, 감별 근거, 동선과 기록.') +
    card(x4[2], 2.08, w4, 3.05, BLUE, '방역관', '번지고 있나?', '집단 이상 탐지, 접촉망, 확산 시나리오, 역학조사 기록.') +
    card(x4[3], 2.08, w4, 3.05, AMBER, '행정', '지역 상황은?', '보급 현황, 성과 지표, 정책 브리핑과 예산 근거.') +
    callout('권한은 서버에서 강제됩니다. 범위 밖 데이터는 감추는 것이 아니라 거부됩니다.'), nxt()),
    '"감춘다"와 "거부한다"의 차이는 조달 심사에서 중요하다. 감추는 것은 화면의 선택이고, 거부하는 것은 구조의 보장이다. '
    '어시스턴트의 장기기억에도 같은 규칙이 적용된다 — 계정 권한 안의 것만 회상한다.')

# ── 45 §09
add(divider(9, '수의 임상 지원',
    '시스템은 후보와 그 근거를 정리합니다. 진단과 처방은 수의사의 권한으로 남습니다.', nxt()),
    '수의사가 있는 자리라면 이 섹션이 신뢰를 좌우한다. 한계부터 말하고 근거 표를 보여준다. '
    '"AI가 진단합니다"라는 말은 절대 하지 않는다.')

# ── 46 감별
add(content('감별', '근거 없는 확률은 쓸모가 없습니다',
    stat(ML, 2.06, 3.33, '6', '사례마다 6개 후보를 순위로 제시하고, 각각의 근거를 함께 보여줍니다') +
    bullets([
        '후보마다 어떤 소견이 지지하고 어떤 소견이 반대하는지 표시한다',
        '자료가 없는 항목은 "없음"으로 표시한다 — 음성으로 세지 않는다',
        '근거가 약한 단일 후보가 높은 확신에 도달하지 못하게 막는다',
        '다음 확인검사와 긴급도를 함께 제안한다',
    ], x=4.48, tx=4.77, w=7.81, y0=2.04) +
    callout('불확실성은 감추지 않고 드러냅니다. 근거가 얇은데 확신해 보이는 숫자는 결함입니다.'), nxt()),
    '고친 결함을 밝힌다 — 초기 버전은 근거가 약한 단일 후보가 확신에 가깝게 표시됐다. '
    '불확실성 prior 보정으로 제거했다. 고친 결함을 말하는 쪽이 어떤 정확도 주장보다 임상 신뢰를 얻는다.')

# ── 47 기록
w3, x3 = row_x(3, gap=0.43)
add(content('기록', '기록이 저절로 남습니다 — 휴약 종료일까지',
    card(x3[0], 2.08, w3, 2.95, TEAL, '처방', '진료 현장에서 바로', '진단·약제·용량·투여경로·기간을 그 자리에서 한 번 입력하면 처방전이 나온다.') +
    card(x3[1], 2.08, w3, 2.95, CORAL, '휴약', '자동으로 계산', '약제에서 종료일이 계산되어 집유·출하 기준으로 표시된다. 암산이 필요 없다.') +
    card(x3[2], 2.08, w3, 2.95, BLUE, '결과', '센서가 확인', '며칠 뒤 시스템이 그 소를 다시 읽어 회복·재발·관찰을 보고한다.') +
    callout('휴약 관리는 식품안전이고, 식품안전은 농장 문제이기 전에 통상 문제입니다.'), nxt()),
    '정부 청중에게는 휴약 카드가 이 섹션에서 가장 센 장이다. 잔류 관리가 지금은 종이와 기억으로 돌아간다. '
    '여기서는 시스템의 속성이 되고, 수출 시장이 요구하는 증빙을 만들어 낸다.')

# ── 48 전문가 루프
add(content('학습', '수의사의 판단이 시스템의 기억이 됩니다', steps([
    ('1', '판단이 말로 나온다', '케이스 검토 중 — "이건 케토시스네"'),
    ('2', '한 번의 확인', '어시스턴트가 한 줄로 묻는다. 동의할 때만 기록된다'),
    ('3', '저장하고 집계한다', '정답 레이블로 남고 알람 유형별로 누적 추적된다'),
    ('4', '다음 판단이 좋아진다', '가이던스는 정해진 규칙으로 다시 만든다 — 모델이 고쳐 쓰지 않는다'),
], y=2.14, h=2.72, sz=13) +
    callout('전문가 판단만 레이블이 됩니다. 농가·행정관의 추측은 학습 데이터 오염을 막기 위해 제외합니다.'), nxt()),
    '4단계는 일부러 보수적으로 만들었다. 모델이 자기 지침을 고쳐 쓰지 않는다. '
    '레이블 통계에서 결정적 규칙으로 가이던스를 만들기 때문에 모든 변경이 검토 가능하다. '
    '"AI가 제멋대로 바뀌는 것을 어떻게 막느냐"에 대한 답이 이것이다.')

# ── 49 §10
add(divider(10, '번식 루프',
    '발정 알람에서 수정 결정으로, 임신감정 결과로, 그리고 다시 다음 추천으로 돌아옵니다.', nxt()),
    '목장이 돈을 가장 빨리 체감하는 지점이다. 공태일은 낙농에서 가장 비싼 숫자이고, '
    '센서와 결정 레이어가 가장 먼저 움직이는 숫자이기도 하다.')

# ── 50 발정
add(content('발정', '추천에는 그 목장의 설정이 실려 있습니다', bullets([
    '수정 적기는 일반값이 아니라 이 목장의 번식 설정으로 계산한다',
    '정액은 근교계수·유전능력·목장의 개량 목표로 제안한다',
    '먼저 사전 점검이 돈다 — 최근 치료, 휴약 상태, 반복 실패 이력',
    '누가 언제 하는지까지 나오고, 수정사에게 그대로 전달된다',
], y0=2.12) + callout('순위는 근교계수·유전능력·목장 목표로 매깁니다. 공급사 가중치는 알고리즘에 존재하지 않습니다.'), nxt()),
    '목장마다 번식 설정이 다르다 — 발정재귀일, 임신감정 시기, 임신기간, 건유 시점. '
    '여기에 일반값을 쓰는 것이 자문 도구가 신뢰를 잃는 가장 흔한 방식이다. '
    '콜아웃은 조달 답변이다. 공급사 편향 질문이 나오면 그대로 읽는다.')

# ── 51 동기화·피드백
add(content('처방과 피드백', '프로토콜은 일정이 되고, 결과는 학습이 됩니다', two_col(
    (BLUE, '발정동기화', '프로토콜이 달력이 된다',
     ['표준 프로토콜 4종을 처방으로 지원한다', '날짜가 찍힌 호르몬 투여·수정 일정으로 펼쳐진다',
      '오늘 할 처치가 작업목록에 나타난다', '승인 규칙은 그대로 적용된다']),
    (GREEN, '닫히는 루프', '결과가 다음 주기를 가르친다',
     ['수정은 정액과 수정사와 함께 기록된다', '임신감정 결과가 정답으로 돌아온다',
      '정액별·개체별 수태율이 추적된다', '그만큼 다음 추천이 조정된다'])) +
    callout('돌아오지 않는 임신감정 결과는 닫히지 않는 루프이고, 나아지지 않는 시스템입니다.'), nxt()),
    '경고할 실패 유형: 수정은 기록하는데 임신감정 결과는 안 넣는 목장. '
    '결과가 없으면 학습이 안 되고 추천 품질이 정체된다. 결과 기록을 1주차부터 일과에 넣는다.')

# ── 52 §11
add(divider(11, '한 농장에서 국가로',
    '한 농장을 돕는 그 신호가, 모이면 국가 방역의 조기 감시망이 됩니다.', nxt()),
    '여기서 프레임을 바꾼다. 지금까지는 한 농장 이야기였다. 여기서부터는 지역과 국가를 생각하게 만든다.')

# ── 53 드릴다운
add(content('드릴다운', '네 단계가 하나의 화면으로 이어집니다', steps([
    ('1', '전국', '시도별 위험등급과 발열률을 한 지도에서'),
    ('2', '시도', '관내 농장의 좌표·두수·알림 건수'),
    ('3', '농장', '우군 상태, 최근 알림, 이동 이력과 접촉'),
    ('4', '개체', '그 소의 체온 곡선과 알람의 해석 근거'),
], y=2.20, h=2.60) + callout('전국 지도에서 소 한 마리의 체온 곡선까지 클릭 네 번. 시스템을 갈아타지 않습니다.'), nxt()),
    '네 개의 화면이 아니라 하나의 연속 동작으로 시연한다. 남겨야 할 인상은 '
    '"전국 숫자와 개별 개체가 같은 데이터"라는 것이지, 대조해야 할 두 개의 보고서가 아니라는 것이다.')

# ── 54 확산
add(content('확산', '접촉망과 확산 모형 — 그림이 아닙니다', two_col(
    (CORAL, '접촉망', '이동 기록에서 만든다',
     ['기록된 이동 이력으로 접촉 개체를 추적한다', '이동 데이터가 없으면 지역 기반으로 대체하고 그렇다고 표시한다',
      '반경 분석으로 거리 안의 농장 수와 두수를 낸다']),
    (BLUE, '확산 시뮬레이션', '미분방정식 모형',
     ['질병별 계수를 넣은 SEIR 모형', '이동제한 시나리오를 수치로 비교한다',
      '인상이 아니라 곡선과 농가 수를 출력한다'])) +
    callout('모형은 결정을 돕습니다. 결정을 대신하지 않고, 특정 발생을 예측하지도 않습니다.'), nxt()),
    '모형이 무엇인지 정확히 말한다. 해당 질병의 공표된 역학 계수를 쓴 것이지 발생 여부를 예측하는 것이 아니다. '
    '여기를 과장하면 방역 당국 앞에서 신뢰를 잃는다.')

# ── 55 신고 이전
add(content('보고 체계', '신고 이전 구간이 이 레이어의 자리입니다', steps([
    ('—', '생체 이상', '몸속에서는 잡히지만 지금 시스템에는 비어 있다'),
    ('—', '임상 증상', '농가 눈에 보이지만 여전히 기록에 없다'),
    ('1', '신고', '국가 시스템이 시작되는 지점'),
    ('2', '확인', '시료 채취와 정밀 검사'),
    ('3', '대응', '이동제한과 방역 조치'),
], y=2.20, h=2.60, sz=12.5, colors=[MUTED, MUTED, TEAL, TEAL, TEAL]) +
    callout('국가 시스템을 대체하지 않습니다. 그 앞의 두 칸을 채워서 넘겨드립니다.'), nxt()),
    '대체하지 않는다는 문장을 일찍, 분명히 말한다. 공식 신고 체계와 경쟁하는 것으로 읽히면 '
    '방역 당국과의 대화는 그 자리에서 끝난다. 우리는 그 체계의 입력이다.')

# ── 56 §12
add(divider(12, '배우되, 책임을 남긴다',
    '네 개의 루프가 시스템을 개선합니다. 넷 다 감사 가능하고, 어느 것도 모델이 자기 규칙을 고쳐 쓰게 두지 않습니다.', nxt()),
    '제목의 두 부분이 똑같이 중요하다. 배우지만 감사할 수 없는 시스템은 규제 산업에 들어갈 수 없다.')

# ── 57 학습 루프
add(content('학습 루프', '네 개의 루프, 넷 다 감사 가능합니다',
    card(x4[0], 2.08, w4, 3.05, CORAL, '전문가 레이블', '임상 정답', '수의사의 판단이 레이블 사례가 된다. 비전문가 추측은 설계상 제외한다.') +
    card(x4[1], 2.08, w4, 3.05, TEAL, '알람 통계', '자기 보정 가이던스', '확진률과 오탐률을 정해진 규칙으로 문장 가이던스로 바꾼다.') +
    card(x4[2], 2.08, w4, 3.05, BLUE, '치료 결과', '실제로 나았는가', '치료 후 센서가 그 소를 다시 읽어 회복과 재발을 보고한다.') +
    card(x4[3], 2.08, w4, 3.05, AMBER, '대화 기억', '지속되는 목장 사실', '사용자가 말한 장비·원칙·제약을 권한 경계 안에서 기억한다.') +
    callout('레이블 10건 미만의 통계는 절대 동작을 바꾸지 않습니다. 작은 표본은 근거가 아닙니다.'), nxt()),
    '10건 기준선을 명시적으로 말할 가치가 있다. 배우는 시스템과 잡음에 과잉반응하는 시스템의 차이가 여기 있다. '
    '감사기관과 연구자 모두 이 임계값을 묻는다.')

# ── 58 거버넌스
add(content('거버넌스', '데이터 주권은 약속이 아니라 구조로 증명합니다', bullets([
    '모든 값은 DB를 거친다 — 모델이 파이프라인을 우회해 스스로 답할 수 없다',
    '모든 도구 호출이 기록된다 — 누가, 언제, 어느 개체를, 무엇을 받았는지',
    '권한 밖 요청은 화면에서 거르는 것이 아니라 게이트웨이에서 거부된다',
    '어시스턴트의 장기기억조차 그 계정의 권한 경계 안에서만 회상된다',
], y0=2.06, pitch=0.64) +
    [sp(ML, 4.76, CW, 0.36, geom='roundRect', fill=WHITE),
     text(1.06, 4.84, CW - 0.62, 0.24,
          '반출 경계 — 원시 센서 데이터, 개체 식별정보, 진료 기록은 국가 인스턴스를 벗어나지 않습니다.',
          12.75, MUTED)] +
    callout('코드는 국경을 넘고, 데이터는 넘지 않습니다.'), nxt()),
    '해외 청중에게는 이 장이 계약을 좌우한다. 주권 질문은 상대가 꺼내기 전에 먼저 말한다 — '
    '요청받고 답하면 방어로 읽히고, 먼저 말하면 설계로 읽힌다. '
    '국가마다 자기 인스턴스를 돌리고, 승인된 익명 집계만 경계를 넘는다.')

# ── 59 §13
add(divider(13, '현장에 심는다',
    '시스템은 아침 일곱 시 반에 실제로 일어나는 일을 바꿀 때에만 값을 합니다.', nxt()),
    '마지막 섹션은 기술이 아니라 습관에 관한 것이다. 실패한 도입의 대부분은 소프트웨어는 잘 돌았고 일과가 없었다.')

# ── 60 리듬
add(content('일과', '아침마다 같은 세 가지', steps([
    ('1', '밤사이 목록을 본다', '아무도 안 보는 동안 무엇이 달라졌는가 — 시간순이 아니라 위험순으로'),
    ('2', '결정 카드를 연다', '그 소에 대해 알아야 할 것이 한 화면에 — 이력, 맥락, 추천'),
    ('3', '한 것을 기록한다', '했다, 또는 안 했다와 이유. 1초 걸리고, 루프를 닫는 유일한 단계다'),
], y=2.20, h=2.60, sz=14) +
    callout('조치가 기록되지 않으면 루프가 닫히지 않고, 시스템은 나아지지 않습니다.'), nxt()),
    '농가에는 이 세 가지만 가르친다. 1주차에 기능을 하나 더 가르칠 때마다 정착률이 떨어진다. '
    '3단계가 사람들이 건너뛰는 단계이고, 제품이 좋아지는지를 결정하는 단계다.')

# ── 61 측정
add(content('측정', '결과만이 아니라 과정을 측정합니다', bullets([
    '조치 기록률 — 알림 중 실제 결정으로 끝난 비율',
    '조기감지 리드타임 — 원래보다 며칠 일찍 확인했는가',
    '공태일과 수태율 — 목장이 체감하는 번식 성적',
    '치료 성공률과 재발 — 조치가 유지됐는지 센서로 확인',
    '두당 일 마진 — 결국 이 모든 것이 도착하는 자리',
], y0=2.06, pitch=0.62, h=0.58) +
    callout('조치 기록률은 나머지 지표가 움직일지를 미리 알려주는 단 하나의 지표입니다.'), nxt()),
    '이 지표들은 파일럿이 끝난 뒤가 아니라 시작 전에 합의한다. 기준선 없는 파일럿은 '
    '아무리 잘 됐어도 90일 뒤에 아무것도 증명하지 못한다.')

# ── 62 30일
add(content('정착', '30일이면 시스템이 습관이 됩니다', steps([
    ('1주', '설치와 연결', '센서 삽입, 기지국 가동, 역할별 계정과 권한 설정'),
    ('2주', '기준선 측정', '지금 목장이 어디에 있는지를 기록한다 — 가장 많이 건너뛰는 단계'),
    ('3주', '3터치 운영', '지원을 붙인 채 일과 운영. 핵심 세 동작만'),
    ('4주', '점검과 조정', '지표를 함께 읽고 알람 설정을 조율하고 일과를 확정한다'),
], y=2.20, h=2.60) +
    callout('2주차를 건너뛰기 쉽습니다. 기준선이 없으면 90일 뒤에 증명할 것이 없습니다.'), nxt()),
    '2주차를 건너뛰자는 요구에 단호하게 맞선다. 모든 목장이 알람을 바로 쓰고 싶어 한다. '
    '기준선 측정 1주가 90일 리뷰를 의견이 아니라 결과로 만든다.')

# ── 63 핵심 정리
def takeaways(page):
    o = [text(ML, 0.60, 5.42, 0.25, '핵심 정리', 10.5, GREEN, b=True),
         text(ML, 1.00, 11.00, 0.95, '이 자리에서 가져갈 세 문장', 34.5, WHITE, b=True)]
    for k, (n, t) in enumerate([
        ('1', '센서는 무언가 달라졌다고 알려주고, 이 레이어는 그래서 무엇을 할지 알려줍니다.'),
        ('2', '알람은 진단이 아니고, AI의 답변은 결정이 아닙니다. 둘 다 사람이 쓰는 재료입니다.'),
        ('3', '조치를 기록해야 루프가 닫히고, 닫히기 전까지는 아무것도 나아지지 않습니다.'),
    ]):
        y = 2.55 + k * 1.16
        o += [sp(ML, y, 0.62, 0.62, geom='ellipse', fill=GREEN),
              text(ML, y + 0.14, 0.62, 0.34, n, 18.75, BG_DARK, b=True, align='ctr'),
              text(1.62, y + 0.06, 10.20, 0.90, t, 20.25, WHITE)]
    return o + chrome(page, dark=True), True
add(takeaways(nxt()),
    '세 번째 문장에서 끝내고 말을 멈춘다. 요약의 요약을 덧붙이지 않는다. '
    '시간이 남으면 "셋 중 우리 목장에서 가장 어려운 게 뭡니까"를 묻는다. 답은 언제나 세 번째다.')

# ══ 1) 원본 35장 번역 ══════════════════════════════════
from pptx import Presentation
from pptx.oxml.ns import qn
import copy

def set_run_font(run):
    rPr = run._r.get_or_add_rPr()
    for tag, face in (('a:latin', 'Calibri'), ('a:ea', KRFONT), ('a:cs', KRFONT)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', face)

def translate(pptx_in, pptx_out):
    prs = Presentation(pptx_in)
    hit = miss = 0
    for si, s in enumerate(prs.slides, 1):
        if si > 35: break
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            for pi, para in enumerate(sh.text_frame.paragraphs):
                runs = para.runs
                if not runs: continue
                cur = ''.join(r.text for r in runs)
                if not cur.strip(): continue
                key = f'{si}.{sh.shape_id}.{pi}'
                new = TR.get(key)
                if new is None and cur.strip() == FOOTER_EN:
                    new = FOOTER_KR
                if new is None:
                    if cur.strip().isdigit() or len(cur.strip()) <= 3:
                        for r in runs: set_run_font(r)
                        continue
                    miss += 1
                    print('  [미번역]', key, cur[:60])
                    for r in runs: set_run_font(r)
                    continue
                runs[0].text = new
                set_run_font(runs[0])
                for r in runs[1:]:
                    r.text = ''
                hit += 1
    prs.save(pptx_out)
    print(f'번역 적용: {hit}건 / 미번역 {miss}건')

# ══ 2) 국문 28장 추가 ══════════════════════════════════
HEAD = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld>'
        '<p:bg><p:bgPr><a:solidFill><a:srgbClr val="{bg}"/></a:solidFill><a:effectLst/></p:bgPr></p:bg>'
        '<p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>')
TAIL = '</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>'
NOTES = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
         '<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
         'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
         'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree>'
         '<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
         '<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/>'
         '<a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>'
         '<p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/>'
         '<p:cNvSpPr><a:spLocks noGrp="1" noRot="1" noChangeAspect="1"/></p:cNvSpPr>'
         '<p:nvPr><p:ph type="sldImg"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp>'
         '<p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/>'
         '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
         '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/>'
         '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r>'
         f'<a:rPr lang="ko-KR" altLang="en-US"><a:ea typeface="{KRFONT}"/></a:rPr>'
         '<a:t>{n}</a:t></a:r></a:p></p:txBody></p:sp></p:spTree></p:cSld>'
         '<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:notes>')

def append_slides():
    base = 35
    pres = open(f'{UN}/ppt/presentation.xml', encoding='utf-8').read()
    prels = open(f'{UN}/ppt/_rels/presentation.xml.rels', encoding='utf-8').read()
    ct = open(f'{UN}/[Content_Types].xml', encoding='utf-8').read()
    sa, ra, ca = [], [], []
    rid, sid = 42, 291
    for k, ((shapes, dark), notes) in enumerate(SLIDES, 1):
        n = base + k
        open(f'{UN}/ppt/slides/slide{n}.xml','w',encoding='utf-8').write(
            HEAD.format(bg=BG_DARK if dark else BG_LIGHT) + ''.join(shapes) + TAIL)
        open(f'{UN}/ppt/slides/_rels/slide{n}.xml.rels','w',encoding='utf-8').write(
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
            f'<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{n}.xml"/>'
            '</Relationships>')
        open(f'{UN}/ppt/notesSlides/notesSlide{n}.xml','w',encoding='utf-8').write(NOTES.format(n=esc(notes)))
        open(f'{UN}/ppt/notesSlides/_rels/notesSlide{n}.xml.rels','w',encoding='utf-8').write(
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster" Target="../notesMasters/notesMaster1.xml"/>'
            f'<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{n}.xml"/>'
            '</Relationships>')
        sa.append(f'<p:sldId id="{sid}" r:id="rId{rid}"/>')
        ra.append(f'<Relationship Id="rId{rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{n}.xml"/>')
        ca.append(f'<Override PartName="/ppt/slides/slide{n}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>')
        ca.append(f'<Override PartName="/ppt/notesSlides/notesSlide{n}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>')
        rid += 1; sid += 1
    open(f'{UN}/ppt/presentation.xml','w',encoding='utf-8').write(pres.replace('</p:sldIdLst>', ''.join(sa)+'</p:sldIdLst>'))
    open(f'{UN}/ppt/_rels/presentation.xml.rels','w',encoding='utf-8').write(prels.replace('</Relationships>', ''.join(ra)+'</Relationships>'))
    open(f'{UN}/[Content_Types].xml','w',encoding='utf-8').write(ct.replace('</Types>', ''.join(ca)+'</Types>'))

if __name__ == '__main__':
    src = sys.argv[1] if len(sys.argv) > 1 else 'src.pptx'
    print('1) 원본 35장 번역')
    translate(src, '_kr_base.pptx')
    print('2) 국문 28장 추가')
    if os.path.isdir(UN): shutil.rmtree(UN)
    zipfile.ZipFile('_kr_base.pptx').extractall(UN)
    append_slides()
    out = 'smaxtec_cowtalk_kr_training_with_ai.pptx'
    if os.path.exists(out): os.remove(out)
    zf = zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED)
    for root, _, files in os.walk(UN):
        for f in files:
            p = os.path.join(root, f); zf.write(p, os.path.relpath(p, UN))
    zf.close()
    os.remove('_kr_base.pptx')
    print(f'완료: {out} — 총 {35+len(SLIDES)}장')
