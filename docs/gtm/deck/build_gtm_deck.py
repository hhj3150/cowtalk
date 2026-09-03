# -*- coding: utf-8 -*-
"""
CowTalk GTM 덱 → 편집 가능한 PPTX 빌더.

단일 소스는 cowtalk-gtm-deck.html 안의 슬라이드 데이터다. 이 스크립트는 그 HTML을
파싱해 python-pptx 네이티브 도형·텍스트로 다시 그린다. 즉 문구를 고칠 곳은 항상
HTML 한 곳이고, PPTX는 파생물이다.

사용:
    pip install python-pptx
    python3 docs/gtm/deck/build_gtm_deck.py [출력경로.pptx]

주의: 한글 폰트는 '맑은 고딕'으로 강제(eastAsia 지정)한다. 웹덱과 글꼴이 다르지만
파워포인트가 설치된 어느 PC에서나 열리는 쪽을 택했다.
"""
import html as htmlmod
import json
import os
import re
import sys
from html.parser import HTMLParser

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches as I, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, 'cowtalk-gtm-deck.html')
OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, 'CowTalk_GTM_Deck.pptx')

# ── 디자인 토큰 (웹덱 라이트 팔레트와 동일) ──────────────────────
INK        = RGBColor(0x0B, 0x22, 0x39)
INK2       = RGBColor(0x24, 0x3B, 0x53)
MUTED      = RGBColor(0x5A, 0x6B, 0x76)
FAINT      = RGBColor(0x8B, 0x9B, 0xA5)
LINE       = RGBColor(0xD6, 0xE0, 0xE4)
LINE_STRONG= RGBColor(0xB6, 0xC6, 0xCD)
TEAL       = RGBColor(0x0E, 0x7C, 0x7B)
MINT       = RGBColor(0x16, 0xB0, 0xA6)
AMBER      = RGBColor(0xB0, 0x6A, 0x16)
AMBER_SOFT = RGBColor(0xF5, 0xEB, 0xDC)
RISK       = RGBColor(0xA8, 0x32, 0x2F)
RISK_SOFT  = RGBColor(0xF7, 0xE7, 0xE6)
SURFACE    = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2   = RGBColor(0xF6, 0xF9, 0xF9)
SURFACE3   = RGBColor(0xEA, 0xF1, 0xF1)

KR   = "맑은 고딕"
MONO = "Consolas"

EW, EH = 13.333, 7.5
ML, MR = 0.78, 0.78
CW = EW - ML - MR
TOP = 0.62


# ── 1. HTML → 슬라이드 데이터 ──────────────────────────────────
class Node:
    __slots__ = ('tag', 'cls', 'kids', 'text')

    def __init__(self, tag='', cls=''):
        self.tag, self.cls, self.kids, self.text = tag, cls, [], ''

    def has(self, c):
        return c in self.cls.split()

    def find_all(self, tag=None, cls=None):
        out = []
        for k in self.kids:
            if isinstance(k, Node):
                if (tag is None or k.tag == tag) and (cls is None or k.has(cls)):
                    out.append(k)
                out.extend(k.find_all(tag, cls))
        return out

    def inner_text(self):
        """마크업을 제거하고 텍스트만 — <strong>은 굵기 마커 \x01로 남긴다."""
        buf = []
        for k in self.kids:
            if isinstance(k, str):
                buf.append(k)
            elif k.tag in ('strong', 'b'):
                buf.append('\x01' + k.inner_text() + '\x01')
            elif k.tag == 'br':
                buf.append('\n')
            else:
                buf.append(k.inner_text())
        return re.sub(r'[ \t]+', ' ', ''.join(buf)).strip()


VOID = {'br', 'img', 'hr', 'input', 'meta', 'link'}


class Tree(HTMLParser):
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.root = Node('root')
        self.stack = [self.root]

    def handle_starttag(self, tag, attrs):
        n = Node(tag, dict(attrs).get('class', '') or '')
        self.stack[-1].kids.append(n)
        if tag not in VOID:
            self.stack.append(n)

    def handle_startendtag(self, tag, attrs):
        self.stack[-1].kids.append(Node(tag, dict(attrs).get('class', '') or ''))

    def handle_endtag(self, tag):
        if tag in VOID:
            return
        for i in range(len(self.stack) - 1, 0, -1):
            if self.stack[i].tag == tag:
                del self.stack[i:]
                return

    def handle_data(self, data):
        self.stack[-1].kids.append(data)


def parse_html(frag):
    t = Tree()
    t.feed(frag)
    return t.root


def js_slides(src_text):
    """HTML 안의 add({...}) 호출을 순서대로 뽑아 dict 리스트로 만든다."""
    slides = []
    for m in re.finditer(r'\nadd\(\{', src_text):
        i = m.end() - 1            # '{' 위치
        depth, j, in_str, quote = 0, i, False, ''
        while j < len(src_text):
            c = src_text[j]
            if in_str:
                if c == '\\':
                    j += 2
                    continue
                if c == quote:
                    in_str = False
            else:
                if c in '`"\'':
                    in_str, quote = True, c
                elif c == '{':
                    depth += 1
                elif c == '}':
                    depth -= 1
                    if depth == 0:
                        break
            j += 1
        slides.append(parse_obj(src_text[i:j + 1]))
    return slides


def parse_obj(block):
    """add({...}) 리터럴에서 필요한 키만 얕게 추출 (중첩 객체 없음 — 안전)."""
    o = {}
    o['t'] = re.findall(r"'([a-z]+)'", re.search(r"t:\s*\[(.*?)\]", block, re.S).group(1))
    o['cover'] = 'cover:true' in block.replace(' ', '')
    for key in ('eyebrow', 'stage', 'title', 'lede'):
        m = re.search(key + r":\s*'((?:[^'\\]|\\.)*)'", block)
        o[key] = m.group(1).replace("\\'", "'") if m else ''
    for key in ('body', 'n'):
        m = re.search(key + r':\s*`(.*?)`(?=\s*,|\s*\}\s*$)', block, re.S)
        o[key] = m.group(1) if m else ''
    m = re.search(r"foot:\s*\['(.*?)','(.*?)'\]", block, re.S)
    o['foot'] = [m.group(1), m.group(2)] if m else None
    return o


# ── 2. 그리기 헬퍼 ────────────────────────────────────────────
def strip_tags(s):
    return re.sub(r'<[^>]+>', '', s)


def rich(s):
    """<em>/<strong> 을 강조 런으로 쪼갠다 → [(text, emphasis)]"""
    parts, pos = [], 0
    for m in re.finditer(r'<(em|strong|b)>(.*?)</\1>', s, re.S):
        if m.start() > pos:
            parts.append((strip_tags(s[pos:m.start()]), False))
        parts.append((strip_tags(m.group(2)), True))
        pos = m.end()
    if pos < len(s):
        parts.append((strip_tags(s[pos:]), False))
    return [(htmlmod.unescape(t), e) for t, e in parts if t]


def _cjk(run):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', run.font.name or KR)


def textbox(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get('align', PP_ALIGN.LEFT)
        if pa.get('sb') is not None:
            p.space_before = Pt(pa['sb'])
        if pa.get('sa') is not None:
            p.space_after = Pt(pa['sa'])
        p.line_spacing = pa.get('ls', 1.24)
        for r in pa['runs']:
            run = p.add_run()
            run.text = r['t']
            f = run.font
            f.name = r.get('f', KR)
            f.size = Pt(r.get('sz', 12))
            f.bold = r.get('b', False)
            f.color.rgb = r.get('c', INK)
            if r.get('sp'):
                run.font._rPr.set('spc', str(int(r['sp'] * 100)))
            _cjk(run)
    return tb


def box(s, x, y, w, h, fill=None, line=None, lw=0.75):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    shp.text_frame.text = ''
    return shp


def hline(s, x, y, w, color=LINE, weight=0.75):
    ln = s.shapes.add_connector(2, I(x), I(y), I(x + w), I(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def label(text, sz=8.5, c=TEAL):
    return {'runs': [{'t': text.upper(), 'sz': sz, 'c': c, 'f': MONO, 'sp': 1.2, 'b': True}]}


def body_runs(text, sz=10.5, c=MUTED):
    """\x01 로 감싼 구간을 굵게."""
    runs, bold = [], False
    for chunk in text.split('\x01'):
        if chunk:
            runs.append({'t': chunk, 'sz': sz, 'c': INK if bold else c, 'b': bold})
        bold = not bold
    return runs or [{'t': '', 'sz': sz, 'c': c}]


# ── 3. 블록 렌더러 ────────────────────────────────────────────
def draw_card(s, node, x, y, w, h):
    fill, edge, kc = SURFACE2, LINE, TEAL
    if node.has('accent'):
        fill, edge = SURFACE3, LINE_STRONG
    if node.has('warn'):
        fill, edge, kc = AMBER_SOFT, AMBER, AMBER
    if node.has('risk'):
        fill, edge, kc = RISK_SOFT, RISK, RISK
    box(s, x, y, w, h, fill, edge)

    paras, pad = [], 0.20
    for k in node.kids:
        if not isinstance(k, Node):
            continue
        if k.has('kicker'):
            paras.append(dict(label(k.inner_text(), 8, kc), sa=4))
        elif k.tag == 'h4':
            paras.append({'runs': [{'t': k.inner_text().replace('\x01', ''), 'sz': 12.5, 'c': INK, 'b': True}], 'sa': 5, 'ls': 1.2})
        elif k.tag == 'p':
            paras.append({'runs': body_runs(k.inner_text(), 9.8), 'sa': 4, 'ls': 1.32})
        elif k.tag == 'ul':
            for li in k.kids:
                if isinstance(li, Node) and li.tag == 'li':
                    paras.append({'runs': [{'t': '· ', 'sz': 9.5, 'c': TEAL, 'b': True}] + body_runs(li.inner_text(), 9.3, INK2), 'sa': 3, 'ls': 1.3})
    if paras:
        textbox(s, x + pad, y + pad, w - pad * 2, h - pad * 2, paras)


def draw_table(s, node, x, y, w, maxh):
    heads = [th.inner_text() for th in node.find_all('th')]
    rows = []
    for tr in node.find_all('tr'):
        tds = [td for td in tr.kids if isinstance(td, Node) and td.tag == 'td']
        if tds:
            rows.append(([td.inner_text() for td in tds],
                         [td.has('num') for td in tds],
                         tr.has('sum')))
    if not heads and not rows:
        return y
    ncol = max(len(heads), max((len(r[0]) for r in rows), default=0))
    numcol = [any(r[1][i] for r in rows if i < len(r[1])) for i in range(ncol)]
    # 첫 열을 넓게, 숫자 열은 좁게
    weights = []
    for i in range(ncol):
        weights.append(0.62 if numcol[i] else (1.55 if i == 0 else 1.25))
    tot = sum(weights)
    widths = [w * v / tot for v in weights]
    xs, acc = [], x
    for cw in widths:
        xs.append(acc)
        acc += cw

    cur = y
    if heads:
        for i, htxt in enumerate(heads[:ncol]):
            al = PP_ALIGN.RIGHT if numcol[i] else PP_ALIGN.LEFT
            textbox(s, xs[i], cur, widths[i] - 0.08, 0.24,
                    [dict(label(htxt, 7.6), align=al)])
        cur += 0.26
        hline(s, x, cur, w, LINE_STRONG, 0.9)
        cur += 0.09

    n = len(rows) or 1
    rowh = min(0.42, max(0.26, (maxh - (cur - y)) / n))
    for cells, nums, is_sum in rows:
        for i, cell in enumerate(cells[:ncol]):
            al = PP_ALIGN.RIGHT if numcol[i] else PP_ALIGN.LEFT
            sz = 9.6
            runs = body_runs(cell, sz, INK if is_sum else INK2)
            if is_sum:
                for r in runs:
                    r['b'] = True
                    r['c'] = INK
            if numcol[i]:
                for r in runs:
                    r['f'] = MONO
            textbox(s, xs[i], cur + 0.04, widths[i] - 0.08, rowh, [dict({'runs': runs}, align=al, ls=1.18)])
        cur += rowh
        hline(s, x, cur, w, LINE_STRONG if is_sum else LINE, 0.9 if is_sum else 0.6)
    return cur


def draw_metrics(s, node, x, y, w, h):
    ms = [k for k in node.kids if isinstance(k, Node) and k.has('metric')]
    if not ms:
        return
    gap = 0.22
    cw = (w - gap * (len(ms) - 1)) / len(ms)
    for i, m in enumerate(ms):
        cx = x + i * (cw + gap)
        col = AMBER if m.has('amber') else TEAL
        hline(s, cx, y, cw, col, 2.0)
        fig = next((k.inner_text() for k in m.kids if isinstance(k, Node) and k.has('fig')), '')
        lab = next((k.inner_text() for k in m.kids if isinstance(k, Node) and k.has('lab')), '')
        fsz = 15 if m.has('wide') else 24
        textbox(s, cx, y + 0.13, cw, 0.5,
                [{'runs': [{'t': fig, 'sz': fsz, 'c': col if m.has('amber') else INK, 'b': True, 'f': MONO}], 'ls': 1.0}])
        textbox(s, cx, y + 0.13 + (0.36 if m.has('wide') else 0.46), cw, h,
                [{'runs': [{'t': lab, 'sz': 9.2, 'c': MUTED}], 'ls': 1.3}])


def draw_flow(s, node, x, y, w, h, vertical=False):
    steps = [k for k in node.kids if isinstance(k, Node) and k.has('step')]
    if not steps:
        return
    if vertical:
        gap = 0.10
        sh = (h - gap * (len(steps) - 1)) / len(steps)
        for i, st in enumerate(steps):
            sy = y + i * (sh + gap)
            box(s, x, sy, w, sh, SURFACE2, LINE)
            sn = next((k.inner_text() for k in st.kids if isinstance(k, Node) and k.has('sn')), '')
            stt = next((k.inner_text() for k in st.kids if isinstance(k, Node) and k.has('st')), '')
            sd = next((k.inner_text() for k in st.kids if isinstance(k, Node) and k.has('sd')), '')
            textbox(s, x + 0.16, sy + 0.10, w - 0.32, sh - 0.2, [
                dict(label(sn, 7.6), sa=2),
                {'runs': [{'t': stt, 'sz': 11.5, 'c': INK, 'b': True}], 'sa': 2},
                {'runs': body_runs(sd, 9.2), 'ls': 1.28},
            ])
        return
    arrow_w = 0.26
    gap = 0.06
    n = len(steps)
    sw = (w - (arrow_w + gap * 2) * (n - 1)) / n
    for i, st in enumerate(steps):
        sx = x + i * (sw + arrow_w + gap * 2)
        box(s, sx, y, sw, h, SURFACE2, LINE)
        sn = next((k.inner_text() for k in st.kids if isinstance(k, Node) and k.has('sn')), '')
        stt = next((k.inner_text() for k in st.kids if isinstance(k, Node) and k.has('st')), '')
        sd = next((k.inner_text() for k in st.kids if isinstance(k, Node) and k.has('sd')), '')
        textbox(s, sx + 0.15, y + 0.14, sw - 0.3, h - 0.28, [
            dict(label(sn, 7.4), sa=3),
            {'runs': [{'t': stt, 'sz': 11.2, 'c': INK, 'b': True}], 'sa': 3, 'ls': 1.18},
            {'runs': body_runs(sd, 9.0), 'ls': 1.3},
        ])
        if i < n - 1:
            textbox(s, sx + sw + gap, y + h / 2 - 0.16, arrow_w, 0.32,
                    [{'runs': [{'t': '→', 'sz': 14, 'c': FAINT}], 'align': PP_ALIGN.CENTER}],
                    anchor=MSO_ANCHOR.MIDDLE)


def draw_layers(s, node, x, y, w, h):
    ls = [k for k in node.kids if isinstance(k, Node) and k.has('layer')]
    if not ls:
        return
    gap = 0.09
    lh = (h - gap * (len(ls) - 1)) / len(ls)
    colors = {'l4': AMBER, 'l3': TEAL, 'l2': MINT, 'l1': INK2}
    for i, ly in enumerate(ls):
        ly_y = y + i * (lh + gap)
        box(s, x, ly_y, w, lh, SURFACE2, LINE)
        col = next((v for k, v in colors.items() if ly.has(k)), TEAL)
        bar = box(s, x, ly_y, 0.045, lh, col, None)
        bar.line.fill.background()
        lno = next((k.inner_text() for k in ly.find_all('div', 'lno')), '')
        name = next((k.inner_text() for k in ly.find_all('div', 'lname')), '')
        owner = next((k.inner_text() for k in ly.find_all('div', 'lowner')), '')
        desc = next((k.inner_text() for k in ly.find_all('div', 'ldesc')), '')
        textbox(s, x + 0.20, ly_y + 0.13, 1.55, lh - 0.2, [dict(label(lno, 7.6, col))])
        textbox(s, x + 1.85, ly_y + 0.10, w - 4.6, lh - 0.16, [
            {'runs': [{'t': name, 'sz': 12, 'c': INK, 'b': True}], 'sa': 2},
            {'runs': body_runs(desc, 9.3), 'ls': 1.28},
        ])
        textbox(s, x + w - 2.55, ly_y + 0.14, 2.35, 0.3,
                [{'runs': [{'t': owner, 'sz': 8.6, 'c': FAINT, 'f': MONO}], 'align': PP_ALIGN.RIGHT}])


def draw_bullets(s, node, x, y, w, h):
    lis = [k for k in node.kids if isinstance(k, Node) and k.tag == 'li']
    sz = 9.6 if node.has('dense') else 10.8
    paras = []
    for li in lis:
        paras.append({'runs': [{'t': '● ', 'sz': sz - 2.5, 'c': TEAL, 'b': True}] + body_runs(li.inner_text(), sz, INK2),
                      'sa': 6, 'ls': 1.34})
    if paras:
        textbox(s, x, y, w, h, paras)


def draw_quote(s, node, x, y, w, h):
    main = ''.join(k if isinstance(k, str) else ('' if (isinstance(k, Node) and k.tag == 'small') else k.inner_text())
                   for k in node.kids)
    main = re.sub(r'\s*\n\s*', '\n', htmlmod.unescape(re.sub(r'<br\s*/?>', '\n', main))).strip()
    smalls = node.find_all('small')
    bar = box(s, x, y, 0.045, h, AMBER, None)
    bar.line.fill.background()
    paras = [{'runs': [{'t': ln, 'sz': 19, 'c': INK}], 'ls': 1.36, 'sa': 2} for ln in main.split('\n') if ln.strip()]
    if smalls:
        paras.append({'runs': [{'t': smalls[0].inner_text().replace('\x01', ''), 'sz': 10, 'c': MUTED}], 'sb': 12, 'ls': 1.3})
    textbox(s, x + 0.32, y, w - 0.32, h, paras, anchor=MSO_ANCHOR.MIDDLE)


def draw_boundary(s, node, x, y, w, h):
    shp = box(s, x, y, w, h, None, LINE_STRONG)
    shp.line.dash_style = 4  # dash
    bt = next((k.inner_text() for k in node.find_all('div', 'bt')), '')
    textbox(s, x + 0.20, y + 0.14, w - 0.4, 0.24, [dict(label(bt, 7.8, FAINT))])
    cards = [k for k in node.kids if isinstance(k, Node) and k.has('card')]
    if cards:
        inner_y = y + 0.48
        inner_h = (h - 0.68 - 0.14 * (len(cards) - 1)) / len(cards)
        for i, c in enumerate(cards):
            draw_card(s, c, x + 0.18, inner_y + i * (inner_h + 0.14), w - 0.36, inner_h)


def render_body(s, body_html, y0, avail_h):
    """s-body 영역을 그린다. 블록을 위에서 아래로 배치."""
    root = parse_html(body_html)
    blocks = [k for k in root.kids if isinstance(k, Node)]
    if not blocks:
        return
    # 높이 배분: cols/layers/flow 는 신축, table/metrics 는 고정 성향
    def weight(b):
        if b.has('metrics'):
            return 0.9
        if b.tag == 'table':
            return 1.4
        return 2.4
    tot = sum(weight(b) for b in blocks)
    gap = 0.20
    usable = avail_h - gap * (len(blocks) - 1)
    y = y0
    for b in blocks:
        h = usable * weight(b) / tot
        draw_block(s, b, ML, y, CW, h)
        y += h + gap


def draw_block(s, b, x, y, w, h):
    if b.has('cols'):
        kids = [k for k in b.kids if isinstance(k, Node)]
        n = len(kids) or 1
        ratios = [1] * n
        if b.has('c25'):
            ratios = [1.25, 1]
        elif b.has('c52'):
            ratios = [1, 1.6]
        ratios = ratios[:n] + [1] * max(0, n - len(ratios))
        tot = sum(ratios)
        gap = 0.22
        acc = x
        for i, k in enumerate(kids):
            cw = (w - gap * (n - 1)) * ratios[i] / tot
            draw_block(s, k, acc, y, cw, h)
            acc += cw + gap
    elif b.has('card'):
        draw_card(s, b, x, y, w, h)
    elif b.tag == 'table':
        draw_table(s, b, x, y, w, h)
    elif b.has('metrics'):
        draw_metrics(s, b, x, y, w, h)
    elif b.has('flow'):
        draw_flow(s, b, x, y, w, h, vertical='flex-direction:column' in (b.cls + str(b.kids)) or _is_vertical(b))
    elif b.has('stack-layers'):
        draw_layers(s, b, x, y, w, h)
    elif b.has('bul'):
        draw_bullets(s, b, x, y, w, h)
    elif b.has('quote'):
        draw_quote(s, b, x, y, w, h)
    elif b.has('boundary'):
        draw_boundary(s, b, x, y, w, h)
    elif b.tag == 'div':
        # 래퍼 div — 자식을 세로로 쌓는다
        kids = [k for k in b.kids if isinstance(k, Node)]
        if not kids:
            return
        gap = 0.18
        each = (h - gap * (len(kids) - 1)) / len(kids)
        yy = y
        for k in kids:
            draw_block(s, k, x, yy, w, each)
            yy += each + gap


_VERTICAL_FLOWS = set()


def _is_vertical(node):
    return id(node) in _VERTICAL_FLOWS


# ── 4. 슬라이드 조립 ──────────────────────────────────────────
def build():
    src = open(SRC, encoding='utf-8').read()
    # 세로 flow 표시(style 인라인) 사전 스캔
    slides = js_slides(src)
    prs = Presentation()
    prs.slide_width, prs.slide_height = I(EW), I(EH)
    blank = prs.slide_layouts[6]

    for idx, d in enumerate(slides, 1):
        s = prs.slides.add_slide(blank)
        accent = AMBER if d['cover'] else TEAL
        rail = box(s, 0, 0, 0.05, EH, accent, None)
        rail.line.fill.background()

        if d['cover']:
            root = parse_html(d['body'])
            mark = next((n.inner_text() for n in root.find_all('div', 'cover-mark')), '')
            h2 = next((n for n in root.find_all('h2', 'cover-h')), None)
            sub = next((n.inner_text() for n in root.find_all('p', 'cover-sub')), '')
            meta = next((n for n in root.find_all('div', 'cover-meta')), None)
            textbox(s, ML, 2.30, CW, 0.3, [dict(label(mark, 10, AMBER))])
            lines = []
            if h2:
                raw = re.sub(r'<br\s*/?>', '\n', ''.join(
                    k if isinstance(k, str) else ('\x02' + k.inner_text() + '\x02') for k in h2.kids))
                for ln in htmlmod.unescape(raw).split('\n'):
                    runs = []
                    for j, part in enumerate(ln.split('\x02')):
                        if part.strip():
                            runs.append({'t': part, 'sz': 40, 'c': TEAL if j % 2 else INK, 'b': True})
                    if runs:
                        lines.append({'runs': runs, 'ls': 1.16, 'sa': 2})
            textbox(s, ML, 2.72, CW, 1.7, lines)
            textbox(s, ML, 4.58, CW * 0.64, 0.9, [{'runs': [{'t': sub, 'sz': 13, 'c': MUTED}], 'ls': 1.45}])
            if meta:
                items = [n.inner_text().replace('\x01', '') for n in meta.kids if isinstance(n, Node)]
                textbox(s, ML, 5.86, CW, 0.4,
                        [{'runs': [{'t': '     '.join(items), 'sz': 9.5, 'c': FAINT, 'f': MONO}]}])
        else:
            textbox(s, ML, TOP, CW * 0.7, 0.26, [dict(label(d['eyebrow'], 9))])
            if d['stage']:
                chip = box(s, EW - MR - 1.62, TOP - 0.06, 1.62, 0.30, None, LINE)
                chip.line.width = Pt(0.75)
                textbox(s, EW - MR - 1.62, TOP - 0.02, 1.62, 0.24,
                        [dict(label(d['stage'], 7.6, FAINT), align=PP_ALIGN.CENTER)])
            trs = []
            for t, em in rich(d['title']):
                trs.append({'t': t, 'sz': 25, 'c': TEAL if em else INK, 'b': True})
            textbox(s, ML, TOP + 0.36, CW * 0.92, 0.9, [{'runs': trs, 'ls': 1.2}])
            ly = TOP + 1.28
            if d['lede']:
                textbox(s, ML, ly, CW * 0.82, 0.6,
                        [{'runs': [{'t': htmlmod.unescape(strip_tags(d['lede'])), 'sz': 11.5, 'c': MUTED}], 'ls': 1.42}])
                ly += 0.62

            foot_h = 0.52 if d['foot'] else 0.0
            render_body(s, d['body'], ly + 0.18, EH - ly - 0.18 - foot_h - 0.62)

            if d['foot']:
                fy = EH - 1.06
                hline(s, ML, fy, CW, LINE, 0.9)
                textbox(s, ML, fy + 0.13, 1.5, 0.26, [dict(label(d['foot'][0], 7.8))])
                ftxt = re.sub(r'<strong>(.*?)</strong>', lambda m: '\x01' + m.group(1) + '\x01', d['foot'][1])
                ftxt = htmlmod.unescape(strip_tags(ftxt).replace('&nbsp;', ' '))
                textbox(s, ML + 1.58, fy + 0.11, CW - 1.58, 0.4,
                        [{'runs': body_runs(ftxt, 10.2, INK2), 'ls': 1.3}])

        textbox(s, EW - MR - 0.6, EH - 0.52, 0.6, 0.26,
                [{'runs': [{'t': f'{idx:02d}', 'sz': 9, 'c': FAINT, 'f': MONO}], 'align': PP_ALIGN.RIGHT}])

        if d.get('n'):
            note = re.sub(r'<[^>]+>', '', d['n'])
            note = htmlmod.unescape(re.sub(r'\n\s+', '\n', note)).strip()
            s.notes_slide.notes_text_frame.text = note

    prs.save(OUT)
    print(f'saved: {OUT}  ({len(slides)} slides)')


if __name__ == '__main__':
    build()
